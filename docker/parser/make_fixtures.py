#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""造 5 份**真**夹具（docx / pptx / pdf / png / doc），给 `scripts/parser.ps1 probe` 喂。

为什么夹具在容器里造：宿主机造不出来 —— python-docx/python-pptx 依赖的 lxml 编译扩展被本机
Smart App Control 拦死，正是本轮要解的问题本体。所以「造夹具」这一步本身就是第一道判据：
它跑通就证明容器里 docx/pptx 的写路径活着，parse 那侧再证明读路径活着。

用法（容器内）：python /app/make_fixtures.py /kbdata/_probe [/hostfonts/simhei.ttf]

字体只在造 png/pdf 夹具时需要（要真中文字形才能测 OCR），**运行时挂宿主机的 C:\\Windows\\Fonts**
即可，镜像里不装中文字体（省 ~250MB，理由见 Dockerfile）。
"""
import os
import subprocess
import sys

TEXT = ['培训报销标准', '境内培训：实报实销，单次上限 3000 元', '境外培训：需总经理审批']


def main(out_dir, font):
    os.makedirs(out_dir, exist_ok=True)
    made = []

    import docx
    d = docx.Document()
    d.add_heading('培训报销制度', level=1)
    d.add_heading('第一节 标准', level=2)
    for t in TEXT[1:]:
        d.add_paragraph(t)
    tb = d.add_table(rows=2, cols=2)
    tb.cell(0, 0).text, tb.cell(0, 1).text = '类别', '上限'
    tb.cell(1, 0).text, tb.cell(1, 1).text = '境内培训', '3000'
    p_docx = os.path.join(out_dir, 'fixture.docx')
    d.save(p_docx)
    made.append(p_docx)

    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = '培训报销标准'
    s.placeholders[1].text = '\n'.join(TEXT[1:])
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = '第二页 审批流程'
    s2.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1)).text_frame.text = '部门经理 → 财务 → 总经理'
    p_pptx = os.path.join(out_dir, 'fixture.pptx')
    prs.save(p_pptx)
    made.append(p_pptx)

    # PDF：带文本层（走 pymupdf4llm 那级），字体嵌进去否则中文是问号
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_font(fontname='CN', fontfile=font)
    for i, t in enumerate(TEXT):
        page.insert_text((72, 100 + i * 30), t, fontname='CN', fontsize=14)
    p_pdf = os.path.join(out_dir, 'fixture.pdf')
    doc.save(p_pdf)
    doc.close()
    made.append(p_pdf)

    # 图片：白底黑字，1000x220。OCR 的判据要的是「中文真被读出来」，不是像素级还原
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new('RGB', (1000, 220), 'white')
    dr = ImageDraw.Draw(img)
    f = ImageFont.truetype(font, 34)
    for i, t in enumerate(TEXT):
        dr.text((20, 20 + i * 60), t, fill='black', font=f)
    p_png = os.path.join(out_dir, 'fixture.png')
    img.save(p_png)
    made.append(p_png)

    # 旧二进制 .doc：用 LibreOffice 从 docx 倒回去（这台机器上没有别的办法造真 .doc；
    # 手写 OLE2 复合文档是自找麻烦）。转出来的是**真** MS Word 97-2003 二进制格式。
    r = subprocess.run(['soffice', '--headless', '--norestore', '--convert-to', 'doc',
                        '--outdir', out_dir, p_docx], capture_output=True, timeout=120)
    p_doc = os.path.join(out_dir, 'fixture.doc')
    if not os.path.isfile(p_doc):
        raise SystemExit(f'造 .doc 夹具失败：{(r.stderr or r.stdout).decode("utf-8", "replace")}')
    made.append(p_doc)

    for p in made:
        print(f'{p}\t{os.path.getsize(p)} bytes')


if __name__ == '__main__':
    main(sys.argv[1] if len(sys.argv) > 1 else '/kbdata/_probe',
         sys.argv[2] if len(sys.argv) > 2 else '/hostfonts/simhei.ttf')
