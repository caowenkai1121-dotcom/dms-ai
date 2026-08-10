"""二进制格式夹具的**配方**（png / pdf 文本层 / pdf 扫描件 / docx / pptx）。

为什么要有这个文件：这五份夹具是不可读的二进制块，直接提交进仓库就是「没人能审、
坏了没人能重造」。文本夹具（md/txt/csv）改一个数字肉眼可见，二进制的不行 ——
判据 token（600/1480/750/3400/5200/2580）到底在不在文件里，只有配方看得出来。
`.xlsx` 那两份历史夹具没有配方，本文件不追溯补（openpyxl 在本机可用，随时能重造）。

用法（本机 = D:\\code\\dms_ai）：
  .venv\\Scripts\\python tools\\kb_fixtures\\make_binary_fixtures.py         # 造本机能造的 + 自检
  .venv\\Scripts\\python tools\\kb_fixtures\\make_binary_fixtures.py --check  # 只自检，不重造
  # 其余四份在容器里造（下面这条原样能跑，实测 python:3.12-slim 拉取 + pip 约 1 分钟）：
  docker run --rm -v "D:\\code\\dms_ai\\tools\\kb_fixtures:/out" \\
    -v "C:\\Windows\\Fonts\\simhei.ttf:/font.ttf:ro" python:3.12-slim \\
    sh -c "pip install -q python-docx python-pptx pillow && python /out/make_binary_fixtures.py container /out"

🔴 **本机只造得出手写 PDF 那一份**，其余四份都得进容器 —— 两个不同的 SAC 拦截，实测：
 ① `lxml` 的编译扩展被拦（`ImportError: DLL load failed while importing etree`）→
    python-docx / python-pptx 全废，**造不了也解析不了** docx/pptx；同一个 DLL 就是
    `/health` 里 `parse_ok.docx=false` 的成因。
 ② Pillow 的 `_imagingft` 也被拦（`DLL load failed while importing _imagingft:
    应用程序控制策略已阻止此文件`）→ `PIL` 本体能 import、`ImageFont.truetype` 直接抛，
    **本机画不出带中文的图**（`load_default()` 只有 ASCII 位图字，对中文 OCR 题没用）。
    所以 png / 扫描件 pdf 也进容器，字体从宿主机 `simhei.ttf` 挂进去（容器 slim 镜像没 CJK 字体）。
手写 PDF 那份是纯 stdlib（zlib 在标准库里），所以不受任何 SAC 拦截影响 —— 这也是它被选中的原因。

旧二进制 Office（.doc / .xls / .ppt）本文件**不造**：那需要 LibreOffice headless
（`soffice --convert-to doc`），本机没有。它们也**刻意没有题** —— 理由见
kb_eval_cases_binary.json 文件头「还缺的判据 ②」：夹具不在盘上时 `missing_fixtures()`
会让整趟变成「执行 0 题」，写了题反而把这个题集也弄哑。
"""
import os, sys, zlib
from pathlib import Path

OUT = Path(__file__).resolve().parent
# 容器里没有 CJK 字体，宿主机的 simhei.ttf 挂进来当 /font.ttf（单面 .ttf，不用给 Pillow 指 index）
FONT = os.environ.get("DMS_FIXTURE_FONT") or (
    "/font.ttf" if Path("/font.ttf").exists() else "C:/Windows/Fonts/simhei.ttf")

# 判据 token 全语料唯一（与 kb_eval_cases.json 的两条写题纪律同口径）。
# 逐个核过：600/1480/750/3400/5200/2580 在 tools/kb_fixtures/ 其余文件里一次都不出现，
# 且**互不为子串** —— `expect.keywords` 是子串比对，2600 会让钉「600」的题假绿，
# 所以 pptx 的预算刻意取 2580 而不是 2600；同理 pptx 里的学时数写成中文「二十六」而不是 26，
# 免得它与 `培训报销_2026新版.md` 的「2026」互相污染（26 是 2026 的子串，当不了判据）。
PNG_LINES = [
    "劳保用品采购与报销标准（扫描通知 · 评测用样本）",
    "全体在职员工适用，试用期员工同等适用。",
    "每人每年采购上限 600 元，凭正式发票申领。",
    "工装破损须交旧领新，遗失自购不予报销。",
]
SCAN_PDF_LINES = [
    "员工食堂就餐补助通知（扫描件 · 评测用样本）",
    "自本月起，在职员工每人每月就餐补助 750 元。",
    "补助计入食堂卡，不折现、不结转下月。",
    "离职当月按实际在职天数折算。",
]
# (字号, 正文)。字号差是给 pymupdf4llm 认标题用的（它按字号推 markdown 的 #），
# 落到 pypdf 那一级只剩逐页纯文本、heading_path 恒空 —— 见题集里 KBB02 的说明。
TEXT_PDF_PAGES = [
    [(17, "员工公寓管理办法（评测用样本，非真实制度）"),
     (11, "本文件是 kb_eval 二进制格式题集的 PDF 文本层语料，改数字会让 KBB02 失败。"),
     (15, "一、申请与分配"),
     (11, "公寓面向异地调入员工，按到岗顺序分配，同批次以家庭人口多者优先。"),
     (11, "单人间与双人间各占一半，不接受指定房号。")],
    [(15, "二、押金与费用"),
     (11, "入住须缴纳押金 1480 元，退房验收无损坏后 30 个自然日内全额退还。"),
     (11, "水电费按表计收，由住户自行承担；宽带由公司统一开通。"),
     (15, "三、退房与验收"),
     (11, "退房须提前 10 个工作日提交申请，行政部到场验收后办理钥匙交接。")],
]

FILES = {
    "png": "劳保用品采购标准_扫描通知.png",
    "pdf_text": "员工公寓管理办法_文本层.pdf",
    "pdf_scan": "食堂就餐补助通知_扫描件无文本层.pdf",
    "docx": "内部推荐奖金办法.docx",
    "pptx": "新员工入职引导.pptx",
}


# ============ Pillow：图片 + 扫描件 PDF ============
def _render(lines, size=(1100, 460), font_px=34):
    """黑字白底、不加噪声不加旋转 —— **刻意的**：OCR 判据要钉住「600」这个数字准不准，
    夹具本身抖一下（噪声/倾斜）判据就变成抛硬币，红了分不清是 OCR 退化还是采样。
    真实扫描件的鲁棒性是另一件事，不该混在同一道题里。"""
    from PIL import Image, ImageDraw, ImageFont
    font = ImageFont.truetype(FONT, font_px)
    img = Image.new("RGB", size, "white")
    d = ImageDraw.Draw(img)
    for i, ln in enumerate(lines):
        d.text((40, 40 + i * int(font_px * 2.1)), ln, fill="black", font=font)
    return img


def make_png(p):
    _render(PNG_LINES).save(p)


def make_pdf_scan(p):
    """Pillow 直接存 PDF＝**只有图、没有文本层**，正是 `_p_pdf` 三级降级后
    `raise ParseError('no_text_layer')` 的那条路（扫描件）。要它出文字只能上 OCR。"""
    _render(SCAN_PDF_LINES).save(p, "PDF", resolution=150.0)


# ============ 手写 PDF 字节流（纯 stdlib）============
# 骨架沿用 tools/parse_probe.py 的 make_pdf（xref 偏移真算，不靠解析器容错重建），
# 但那份正文是 ASCII —— 它注释里那句「base-14 字体下中文抽不出来」正是这里要绕的：
# 检索是中文的，抽不出中文这份夹具等于什么都没测。
#
# 绕法：Type0 字体（BaseFont STSong-Light，**不嵌字体**）+ 码位＝UTF-16BE + 自写 `/ToUnicode`。
# 抽取靠 ToUnicode；渲染靠 `/UniGB-UCS2-H` 这个**预定义 CMap**（UCS-2 → Adobe-GB1 字形号，
# 阅读器用自带 CJK 替代字体就能显示），于是既抽得对也大致看得见，
# 还省掉一个 TTF 子集器和 15MB 的 simsun.ttc。
# ponytail: 不嵌字体 → 渲染依赖阅读器的 CJK 替代字。肉眼效果不满意就等容器里有 LibreOffice 后
# 用 `soffice --convert-to pdf` 重造这份夹具（判据 token 不变即可）。
#
# 实测（本机 .venv，pypdf 6.14.2，两种 Encoding 各生成一份比对）：
# `/Identity-H` 与 `/UniGB-UCS2-H` **两者 extract_text() 都原样出中文** ——
# pypdf 认的是 ToUnicode，不看 Encoding 名。既然抽取上等价，就选能渲染的那个。
# ToUnicode 是不能省的那一半：没有它，Identity-H 的 CID 无处可映。
PDF_ENC = b"/UniGB-UCS2-H"


def _tounicode_cmap(chars):
    out = ["/CIDInit /ProcSet findresource begin\n12 dict begin\nbegincmap\n"
           "/CIDSystemInfo <</Registry (Adobe) /Ordering (UCS) /Supplement 0>> def\n"
           "/CMapName /Adobe-Identity-UCS def\n/CMapType 2 def\n"
           "1 begincodespacerange\n<0000> <FFFF>\nendcodespacerange\n"]
    cs = sorted(chars)
    for i in range(0, len(cs), 100):        # bfchar 每块上限 100 条（PDF 规范）
        g = cs[i:i + 100]
        out.append(f"{len(g)} beginbfchar\n")
        out += [f"<{ord(c):04X}> <{ord(c):04X}>\n" for c in g]
        out.append("endbfchar\n")
    out.append("endcmap\nCMapName currentdict /CMap defineresource pop\nend\nend\n")
    return "".join(out).encode("latin-1")


def _page_content(lines):
    out = ["BT\n"]
    y = 780
    for size, text in lines:
        hexs = text.encode("utf-16-be").hex().upper()
        out.append(f"/F1 {size} Tf 1 0 0 1 56 {y} Tm <{hexs}> Tj\n")
        y -= int(size * 2.0)
    out.append("ET\n")
    return "".join(out).encode("latin-1")


def make_pdf_text(p, pages=None):
    pages = pages or TEXT_PDF_PAGES
    chars = {c for pg in pages for _, t in pg for c in t}
    cmap = _tounicode_cmap(chars)
    n_page_objs = len(pages) * 2                       # 每页：页对象 + 内容流
    kids = " ".join(f"{7 + 2 * i} 0 R" for i in range(len(pages)))
    objs = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[%s]/Count %d>>" % (kids.encode(), len(pages)),
        b"<</Type/Font/Subtype/Type0/BaseFont/STSong-Light/Encoding" + PDF_ENC
        + b"/DescendantFonts[4 0 R]/ToUnicode 6 0 R>>",
        b"<</Type/Font/Subtype/CIDFontType0/BaseFont/STSong-Light/CIDSystemInfo"
        b"<</Registry(Adobe)/Ordering(GB1)/Supplement 2>>/FontDescriptor 5 0 R/DW 1000>>",
        b"<</Type/FontDescriptor/FontName/STSong-Light/Flags 4/FontBBox[-25 -254 1000 880]"
        b"/ItalicAngle 0/Ascent 880/Descent -254/CapHeight 880/StemV 80>>",
        b"<</Length %d>>stream\n" % len(cmap) + cmap + b"\nendstream",
    ]
    assert len(objs) == 6, "页对象从 7 号起，改了上面的对象数就要同步 kids 与 Contents 的编号"
    for i, pg in enumerate(pages):
        body = zlib.compress(_page_content(pg))        # 压缩：证明解析器真在解流，不是 grep 明文
        objs.append(b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 595 842]/Contents %d 0 R"
                    b"/Resources<</Font<</F1 3 0 R>>>>>>" % (8 + 2 * i))
        objs.append(b"<</Length %d/Filter/FlateDecode>>stream\n" % len(body) + body + b"\nendstream")
    assert len(objs) == 6 + n_page_objs
    buf, offsets = bytearray(b"%PDF-1.4\n"), []
    for i, o in enumerate(objs, 1):
        offsets.append(len(buf))
        buf += b"%d 0 obj" % i + o + b"endobj\n"
    xref_at = len(buf)
    buf += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objs) + 1)
    for off in offsets:
        buf += b"%010d 00000 n \n" % off
    buf += b"trailer<</Size %d/Root 1 0 R>>\nstartxref\n%d\n%%%%EOF\n" % (len(objs) + 1, xref_at)
    p.write_bytes(bytes(buf))


# ============ 容器里跑的那半：docx / pptx ============
def make_docx(p):
    """判据是两个数：3400 在段落里、5200 只在**表格**里。
    钉 5200 是因为 `_p_docx` 的 `tbl` 分支（表格 → `a | b | c` 行）是条独立代码路径，
    只钉段落的话表格整体丢掉也照样绿 —— 而丢的是一张标准表。"""
    import docx
    d = docx.Document()
    d.add_heading("内部推荐奖金办法（评测用样本，非真实制度）", 1)
    d.add_paragraph("本文件是 kb_eval 二进制格式题集的 docx 语料，改数字会让 KBB03 失败。")
    d.add_heading("一、奖金标准", 2)
    d.add_paragraph("推荐人在被推荐人入职满三个月后，可获一次性推荐奖金 3400 元。")
    d.add_paragraph("同一候选人被重复推荐时不重复计奖，以最早提交的推荐记录为准。")
    d.add_heading("二、按职级发放", 2)
    t = d.add_table(rows=0, cols=3)
    for row in [("职级", "奖金（元）", "发放时点"), ("普通岗", "3400", "入职满三个月"),
                ("主管岗", "5200", "入职满六个月")]:
        cells = t.add_row().cells
        for c, v in zip(cells, row):
            c.text = v
    d.add_paragraph("被推荐人入职未满六个月离职的，已发奖金全额扣回。")
    d.save(p)


def make_pptx(p):
    import pptx
    prs = pptx.Presentation()
    for title, body in [
        ("新员工入职引导（评测用样本）", ["本 PPT 是 kb_eval 二进制格式题集的 pptx 语料。"]),
        ("培训预算与课时", ["每位新员工人均培训预算 2580 元。",
                            "引导课时合计二十六学时，入职两周内完成。"]),
        ("考核与转正", ["引导期结束进行一次考核，未通过可补考一次。"]),
    ]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        s.placeholders[1].text_frame.text = "\n".join(body)
    prs.save(p)


# ============ 自检：造完必须证明判据 token 真在文件里 ============
def check():
    """每份夹具**至少证明一件事**。没有这一步，"造好了" 只等于 "文件存在"。"""
    bad = []
    for k, name in FILES.items():
        if not (OUT / name).exists():
            bad.append(f"{name} 不存在（{'容器里造' if k in ('docx', 'pptx') else '本机可造'}）")
    p = OUT / FILES["pdf_text"]
    if p.exists():
        from pypdf import PdfReader
        r = PdfReader(p)
        txt = "".join(pg.extract_text() or "" for pg in r.pages)
        if len(r.pages) != 2:
            bad.append(f"{p.name}: {len(r.pages)} 页（应 2 页，页归属是 page 字段的来源）")
        for w in ("1480", "押金", "员工公寓"):
            if w not in txt:
                bad.append(f"{p.name}: 抽取的文本层缺[{w}] → 判据 token 进不了索引")
        print(f"  pdf 文本层：{len(r.pages)} 页 {len(txt)} 字，1480 {'在' if '1480' in txt else '不在'}")
    p = OUT / FILES["pdf_scan"]
    if p.exists():
        from pypdf import PdfReader
        txt = "".join(pg.extract_text() or "" for pg in PdfReader(p).pages).strip()
        # 反向：这份**必须**抽不出文字，否则它就不是扫描件夹具、KBB05 测不到 OCR
        if txt:
            bad.append(f"{p.name}: 竟抽出了 {len(txt)} 字文本 → 它不是无文本层的扫描件")
        print(f"  pdf 扫描件：文本层 {len(txt)} 字（应为 0）")
    p = OUT / FILES["png"]
    if p.exists():
        from PIL import Image
        with Image.open(p) as im:
            print(f"  png：{im.size} {im.mode}")
            if im.size[0] < 600:
                bad.append(f"{p.name}: 宽 {im.size[0]}px 太小，OCR 认不准数字")
    for k in ("docx", "pptx"):
        p = OUT / FILES[k]
        if p.exists():
            import zipfile
            with zipfile.ZipFile(p) as z:                 # 不 import lxml：本机它是坏的
                parts = [n for n in z.namelist() if n.endswith(".xml")]
                blob = b"".join(z.read(n) for n in parts)
            need = {"docx": ("3400", "5200"), "pptx": ("2580",)}[k]
            missing = [w for w in need if w.encode() not in blob]
            bad += [f"{p.name}: 包内 XML 缺判据 token[{w}]" for w in missing]
            print(f"  {k}：{len(parts)} 个 XML 部件，判据 token "
                  f"{'全在' if not missing else '缺 ' + ','.join(missing)}")
    for b in bad:
        print(f"  ✗ {b}")
    return bad


if __name__ == "__main__":
    args = sys.argv[1:]
    if args and args[0] == "container":        # 被 SAC 挡住的四份（docx/pptx 要 lxml，png/扫描件要 _imagingft）
        OUT = Path(args[1]) if len(args) > 1 else OUT
        make_docx(OUT / FILES["docx"])
        make_pptx(OUT / FILES["pptx"])
        make_png(OUT / FILES["png"])
        make_pdf_scan(OUT / FILES["pdf_scan"])
        print(f"✅ 容器内造好 {FILES['docx']} / {FILES['pptx']} / {FILES['png']} / {FILES['pdf_scan']}")
        sys.exit(0)
    if "--check" not in args:
        make_pdf_text(OUT / FILES["pdf_text"])
        print(f"✅ 本机造好 {FILES['pdf_text']}（纯 stdlib）；另四份见文件头的 docker 一行")
    bad = check()
    sys.exit(1 if bad else 0)
