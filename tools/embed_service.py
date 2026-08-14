#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""向量层 + 文档服务：bge-small-zh-v1.5(512维, fastembed 本地自包含)。
  build    —— 为 meta.table_doc / sql_exemplar / element / datasource 算 embedding 存 pgvector
               + HNSW 索引（离线跑；`--ds` 限定注册表行属于哪个源，默认 dms）
  revec    —— build 的第五个目标，**单独跑**（`revec` 或 `build --revec`）：
               补 kb.chunk 里 embedding IS NULL 的行 + 把 kb.doc 从 chunked 推到 embedded。
               仍缺一行即非 0 退出（理由见 `revec_exit`）
  serve    —— 常驻 HTTP :8077（裁决 V1：文档解析复用本进程本端口）
               POST /embed  {"texts":[...],"query":true} → {"embeddings":[[...]]}
               POST /parse  {"path":"<绝对路径>","mime":""} → {"blocks":[...],"page_count":n,"sheets":[...]}
               POST /chunk  {"blocks":[...],"target_tokens":400,"overlap":60} → {"chunks":[...]}
               GET  /health → {"ok","model","dim","parse_ok":{...},"parse_caps":{".docx":{ok,why}, ".pdf":{...,tiers:{text,ocr}}}}
  selftest —— 自造 md/csv/json/html 跑一遍 parse+chunk（不需要任何第三方解析库）+ 把**每种扩展名**的
               可用/不可用与原因列全 + 钉住能力表纪律（见 `_selftest_caps`）
               + 扫描件 PDF 的 OCR 档判定与夹具（见 `_selftest_pdf_scan`）
用法: python embed_service.py build [--ds dms] | revec | serve [port] [host] | selftest

外部依赖的位置可用环境变量覆盖（容器里可能不在 PATH 上）：
  DMS_SOFFICE   LibreOffice headless 可执行文件（旧二进制 .doc/.xls/.ppt 靠它转格式）
  DMS_TESSERACT tesseract 可执行文件（图片 OCR）
  DMS_OCR_LANG  OCR 语言包，默认 chi_sim+eng
"""
import os, sys, json, re, csv, io, math, itertools, importlib.util, shutil, subprocess, tempfile, threading, urllib.request
from html.parser import HTMLParser   # _p_html 去标签用标准库（SAC 拦的是编译扩展，stdlib 两侧都一定有）
# pythonw（stdout=None）/pytest 捕获流没有 reconfigure：缺席就跳过，不许为编码起不来
for _s in (sys.stdout, sys.stderr):
    if getattr(_s, 'reconfigure', None):
        _s.reconfigure(encoding='utf-8')

MODEL = 'BAAI/bge-small-zh-v1.5'
DIM = 512
QUERY_INSTRUCT = '为这个句子生成表示以用于检索相关文章：'


def pg_conf():
    """自有 PG 的连接参数，来自 `settings.json` 的 `pg_url`。

    🔴 这里原先写着**明文口令**（原文已删），违反「明文凭据只在 settings.json」。
    而 settings 里本来就有 `pg_url` —— 也就是说那是无谓的第二份真相源：
    改了 settings 的口令，本服务照旧拿旧口令连，连不上还得去猜为什么。
    **惰性读**（不是模块顶层常量）：`serve`/`selftest` 这些不碰 PG 的路径，
    缺 settings.json 时不该起不来 —— 只有真去连库的那一步才需要凭据。"""
    from settings import pg_kwargs   # tools/ 与本文件同目录，走 sys.path[0]
    return pg_kwargs()

_embedder = None
def embedder():
    global _embedder
    if _embedder is None:
        from fastembed import TextEmbedding
        _embedder = TextEmbedding(model_name=MODEL)
    return _embedder

# 推理的全局锁：fastembed/onnxruntime 的 session 不是线程安全的，所以推理**仍然串行**。
# 它存在的唯一理由是让 `serve` 能换成 ThreadingHTTPServer（见那里的注释）——
# /parse（上限 120s）与 /health 从此不再排在一次 275 块的 /embed 后面。
# 实测（本文件 `_selftest_serve_unblocked`；数字出自旧版 0.6s 桩，现桩睡 2s，量级结论不变）：
# 单线程 /health 要 0.605s，多线程 0.002s。顺带把首次惰性加载也串了：两个线程同时进来会各造一个 TextEmbedding。
# ponytail: 一把进程级锁 —— 本进程只有一个模型；真要并发推理得起多进程，那时再拆。
_EMBED_LOCK = threading.Lock()

def embed(texts, is_query=False):
    if is_query:
        texts = [QUERY_INSTRUCT + t for t in texts]
    with _EMBED_LOCK:
        return [v.tolist() for v in embedder().embed(texts)]

# ============ 文档服务：解析（K1）============
# 解析库一律惰性 import：缺依赖只让该类型报 unsupported，embed 功能不受影响。
MAX_ROWS, MAX_COLS = 200000, 200   # 单 sheet 上限，与 knowledge::tabular 一致；超出即报错不截断

class ParseError(Exception):
    """固定 JSON 形状 + HTTP 状态的解析失败（error 码给 Rust 侧枚举）"""
    def __init__(self, error, detail='', status=422):
        super().__init__(error)
        self.status = status
        self.payload = {'error': error}
        if detail:
            self.payload['detail'] = detail

def _push(stack, level, title):
    """标题栈：同级或更深的先弹出，再压入当前标题"""
    while stack and stack[-1][0] >= level:
        stack.pop()
    stack.append((level, title))

def _path(stack):
    return ' > '.join(t for _, t in stack)

def _blk(text, page, stack):
    return {'text': text, 'page': page, 'heading_path': _path(stack)}

# 中文文档常写 `#一级标题`（# 后不空格）：放行，但标题首字符须非 ASCII —— `#tag`/`#123` 不算标题
_H_MD = re.compile(r'^(#{1,6})(?:\s+|(?=[^\x00-\x7f]))(.+?)\s*#*$')

def md_blocks(text, page, stack):
    """markdown/纯文本 → blocks：# 层级维护 heading_path，空行分段。pdf 的 markdown 也走这里"""
    out, buf = [], []
    def flush():
        t = '\n'.join(buf).strip()
        buf.clear()
        if t:
            out.append(_blk(t, page, stack))
    for line in text.splitlines():
        m = _H_MD.match(line.strip())
        if m:
            flush()
            _push(stack, len(m.group(1)), m.group(2).strip())
            out.append(_blk(m.group(2).strip(), page, stack))
        elif line.strip():
            buf.append(line.rstrip())
        else:
            flush()
    flush()
    return out

def _read_text(path):
    """中文文档常是 GBK：utf-8 → gbk → 替换式兜底，绝不因编码整份失败。
    Windows 导出的 UTF-16（带 BOM）先拦：落入 gbk 会解出夹 NUL 的乱码而不是正确解码。"""
    with open(path, 'rb') as f:
        raw = f.read()
    if raw[:2] in (b'\xff\xfe', b'\xfe\xff'):
        return raw.decode('utf-16')      # BOM 在，utf-16 解码器自己分大小端
    for enc in ('utf-8-sig', 'gbk'):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            pass
    return raw.decode('utf-8', 'replace')

def _p_text(path):
    # 滤掉没有任何词字符的块（md 的 `---`/`***` 分隔线），与 `_p_pdf` 的过滤同一口径
    return [b for b in md_blocks(_read_text(path), None, [])
            if re.search(r'\w', b['text'])], 0, []

def _p_json(path):
    """JSON → 美化后的 ```json 代码块（口径：json 转代码块，与文本同一条分块链）。
    校验失败不拒收 —— 按原文入库仍检索得到，但必须在 notes 留痕：
    静默降级成纯文本而不说，界面会以为结构化解析成功。"""
    text = _read_text(path)
    try:
        pretty = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
    except ValueError as e:
        return md_blocks(text, None, []), 0, [], [f'JSON 校验失败（{e}），已按纯文本入库']
    return md_blocks(f'```json\n{pretty}\n```', None, []), 0, []


class _HtmlToText(HTMLParser):
    """HTML → 纯文本：script/style/noscript 整棵丢弃，块级标签换成换行，实体反转义。
    只出文本 —— 任何标签结构都不进检索正文（预览侧原件也只按 text/plain 展示）。"""
    _BLOCK = {'p', 'div', 'br', 'hr', 'li', 'tr', 'table', 'ul', 'ol', 'pre', 'title',
              'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'section', 'article', 'header', 'footer', 'blockquote'}
    _SKIP = {'script', 'style', 'noscript'}

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.parts, self.skip = [], 0

    def handle_starttag(self, tag, attrs):
        if tag in self._SKIP:
            self.skip += 1
        elif tag in self._BLOCK:
            self.parts.append('\n')

    def handle_endtag(self, tag):
        if tag in self._SKIP:
            self.skip = max(0, self.skip - 1)
        elif tag in self._BLOCK:
            self.parts.append('\n')

    def handle_data(self, data):
        if not self.skip:
            self.parts.append(data)

    def text(self):
        raw = re.sub(r'\n+', '\n', ''.join(self.parts))
        return '\n'.join(x for x in (' '.join(l.split()) for l in raw.split('\n')) if x)


def _p_html(path):
    """HTML 去标签按纯文本入库，只用标准库。剥完无文本不报错：
    空 blocks 由 Rust 入库链统一判「文档里没有可索引的文本」（与 _p_text 同口径）。"""
    parser = _HtmlToText()
    parser.feed(_read_text(path))
    parser.close()
    # 同 `_p_text`：无词字符块（剥完标签只剩分隔线的段）不进 blocks
    return [b for b in md_blocks(parser.text(), None, [])
            if re.search(r'\w', b['text'])], 0, []

def _p_pdf(path):
    """三级降级：pymupdf4llm（**保章节层级**）→ fitz（逐页纯文本）→ pypdf（逐页纯文本）。

    🔴 **许可分层是刻意的**：前两个都是 **AGPL-3.0**（pymupdf4llm 与 PyMuPDF 同源）。
    业主已裁决 AGPL-3.0 可用，所以第一级现在是**推荐**的那一级 —— 只有它给出章节层级：
    它输出 markdown，`#` 标题经 `md_blocks` 变成块的 `heading_path`（与 docx 的 Heading 同一条通道，
    检索引用因此能说清「在第几章」）。第二三级只出逐页纯文本，`heading_path` 恒为空串。
    第三级 `pypdf` 是 **BSD-3-Clause**、纯 Python 无二进制依赖，留作兜底：
    法务口径变了、或本机 Smart App Control 拦掉编译产物时，PDF 仍能进库。
    本机 `.venv` 实测正是这一档：pymupdf4llm/fitz 都没装上，`parse_ok['pdf']=True`
    而 `parse_caps()['.pdf']['why']` 报「pypdf 兜底（BSD-3）：只有逐页纯文本，heading_path 为空」。
    装到哪一级必须照实上报 —— 不上报，运维就无从知道自己拿到的是「带章节」还是「纯文本」
    那种质量（两者 `ok` 都是 true，肉眼分不出来）。

    降级用的 `except Exception` 不是偷懒：本机实测 **DLL 加载失败是 OSError 不是 ImportError**
    （`DLL load failed while importing etree`），只捕 ImportError 会让 PyMuPDF 装了但加载不了时
    整份 PDF 报 500，而不是降到 pypdf。同 `_have` 的口径。

    Y1 高阶档：**低文本量即判扫描件**（`_pdf_is_scanned`），低文本页经 `_pdf_ocr_fill` 渲染成图
    走 OCR（千问 vision 优先、tesseract 降级，与 `_p_image` 同一通道）。三级文本引擎共用这一档。
    """
    try:
        import pymupdf4llm
    except Exception:
        return _pdf_fitz(path)
    try:
        pages = pymupdf4llm.to_markdown(path, page_chunks=True)
    except Exception:
        # 运行期异常（损坏/异形 PDF）也续降一级：import 能成不代表转换能成，不许整份 500
        return _pdf_fitz(path)
    stack, out, page_chars = [], [], {}
    for i, pg in enumerate(pages, 1):
        no = (pg.get('metadata') or {}).get('page', i)
        raw = pg.get('text') or ''
        page_chars[no] = _page_chars(raw)   # 扫描件判定用**过滤前**的原始页文本；键与块的页号同源
        out += md_blocks(raw, no, stack)
    # 🔴 **先滤掉没有任何词字符的块，再判空。** 缺这一行的后果是实测出来的：
    # `pymupdf4llm.to_markdown` 每页尾部会输出一条 markdown 分隔线 `-----`，
    # `md_blocks` 把它当正文块 emit，于是 `if not out` 对**扫描件恒不成立** ——
    # 同一份无文本层的 PDF：容器（pymupdf4llm 一级）返 `200 {"blocks":[{"text":"-----"}]}`，
    # 宿主机（pypdf 三级）返 `422 no_text_layer`。
    # 前者意味着 Rust 侧拿到 1 块 → `status=embedded` → 界面「已入库」→ 问什么都答不出来，
    # 而那条 `-----` 还会进向量索引、甚至成为引用。正是本仓反复抓的静默失败族。
    # 第二个面（同一根因）：有文本层的 PDF 每页也会多一个 `-----` 垃圾块，
    # 被 `chunk_blocks` 按 heading_path 合并时拼进相邻 chunk 的正文，污染入库原文与原文回查判据。
    # 实测：员工公寓管理办法_文本层.pdf（2 页）从 12 块降到 10 块，第 5/11 块那两条 `-----` 消失。
    out = [b for b in out if re.search(r'\w', b.get('text', ''))]
    # 低文本量即判扫描件（Y1）：低文本页渲染成图走 OCR —— 判定/页数护栏/失败口径全在
    # `_pdf_ocr_fill`（钉在 `_selftest_pdf_scan`），实测动机写在那里，不重复。
    notes = _pdf_ocr_fill(path, len(pages), page_chars, out)
    # 第三位恒是 sheets（PDF 没有工作表 → 空），说明走第四位 notes（理由见 parse_doc）
    return out, len(pages), [], notes


def _pdf_page_ocr(path, page_no):
    """把 PDF 某一页渲染成图再 OCR（`fitz` + 已有的 `_p_image` 通道）。

    只在 `_pdf_ocr_fill` 补页时调用，**不登记进 `CAPS`**（`.pdf` 的实现者只有 `_p_pdf` 一个）。
    dpi 取 `OCR_DPI`（默认 200，实测折中：150 下 12pt 中文 tesseract 常丢字，300 单页多花 ~1.5 倍）。
    单页成本 ≈ fitz 渲染 0.1s + 识别 0.2~1s；「页数 × 单页成本」就是 `OCR_PAGE_CAP` 护栏的理由。
    """
    import fitz
    with fitz.open(path) as doc:
        pix = doc[page_no - 1].get_pixmap(dpi=OCR_DPI)
        png = pix.tobytes('png')
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
        f.write(png)
        tmp = f.name
    try:
        # ⚠️ 只取第一位。`_p_image` 返的是**4 元**（blocks/frames/sheets/notes）——
        # 写成 `blocks, n, skipped = _p_image(tmp)` 会 `too many values to unpack`，
        # 而它发生在 `_p_pdf` 内部、表现成整份 PDF **HTTP 500**（实测踩过）。
        blocks = _p_image(tmp)[0]
    finally:
        os.unlink(tmp)
    # `_p_image` 按帧编号（渲染出来的单页图恒 1 帧），改回真实页号；heading 标明是 OCR 来的
    for b in blocks:
        b['page'] = page_no
        b['heading_path'] = f'第 {page_no} 页（OCR）'
    return blocks


def _page_chars(text):
    """页文本量口径：非空白字符数。比「词字符」宽（中文标点也是内容）；
    pymupdf4llm 每页尾部分隔线 `-----` 贡献 5 个字符，对 50 的阈值没有影响。"""
    return len(re.sub(r'\s+', '', text or ''))


def _pdf_is_scanned(total_chars, page_count):
    """整份判扫描件（Y1 裁决阈值）：全文 < PDF_DOC_MIN_CHARS，或页均 < PDF_PAGE_MIN_CHARS。
    0 页的退化输入不算扫描件（交给 `_pdf_ocr_fill` 的 `not out` 统一报 no_text_layer）。"""
    if page_count <= 0:
        return False
    return total_chars < PDF_DOC_MIN_CHARS or total_chars / page_count < PDF_PAGE_MIN_CHARS


def _pdf_low_text_pages(page_chars, page_count):
    """要 OCR 补的页号：该页文本层 < PDF_PAGE_MIN_CHARS 个非空白字符（缺席页按 0 计）。
    原「没有任何块的页」判据是它的特例（0 < 50 恒真）—— 低文本页（垃圾文本层）是新抓的形态。"""
    return [i for i in range(1, page_count + 1) if page_chars.get(i, 0) < PDF_PAGE_MIN_CHARS]


def _pdf_ocr_fill(path, page_count, page_chars, out, ocr_fn=None, cap=None):
    """低文本页渲染成图走 OCR，把块补进 `out`（原地），返回 notes（人话，进响应第四位）。

    `ocr_fn(path, page_no) -> blocks`，默认 `_pdf_page_ocr`；参数化只为 selftest 用桩驱动。
    判定与失败口径（`_selftest_pdf_scan` 逐条钉住）：

    🔴 为什么逐页都要补（只判整份是不够的）：实测造一份 2 页 PDF（页 1 真文本层、页 2 只有
    一张图）→ `200 page_count=2`，页 2 图上的字一个都没进索引，文档照旧推 `embedded`。
    带扫描附件/签字页的 PDF 在制度库里非常常见 —— 那是「静默成功」，比整份失败坏得多
    （整份失败会 422，用户知道要先 OCR；丢半份不会报错，用户以为全文都在里面）。

      - 整份判扫描件（`_pdf_is_scanned`）时 **need = 全部页**：「每页 60 个垃圾字符」那种
        文本层单页看能过 50 的线，合计仍是什么都没有 —— 判了扫描件，每页都可疑。
      - 需要 OCR 的页 > cap → `too_large` 响亮失败，**不「OCR 前 N 页然后报已入库」**。
      - 单页 OCR 失败（含引擎缺席/意外异常）不炸整份：记入 notes（会经 Rust `notes` 带到界面）。
      - OCR 成功的页：块按页号排回（chunk 的相邻关系/span 依赖页序）。文本层那 <50 字符的
        残块**保留**：页脚/页码在 OCR 文本里通常也读得出，重复上限 49 字符/页，可接受；
        丢掉才是把「唯一的真文本」赌在 OCR 质量上。
      - 整份判扫描件、OCR 一字未补、文本层合计仍 < PDF_DOC_MIN_CHARS → `no_text_layer`：
        带着零星字符按「已入库」静默过去，正是本仓反复抓的那个失败族。
        （页均腿判扫描但全文 ≥200 的，保留文本层 + 留痕，不升级为整份失败。）
    """
    ocr_fn = ocr_fn or _pdf_page_ocr
    cap = OCR_PAGE_CAP if cap is None else cap
    total = sum(page_chars.get(i, 0) for i in range(1, page_count + 1))
    scanned = _pdf_is_scanned(total, page_count)
    need = list(range(1, page_count + 1)) if scanned else _pdf_low_text_pages(page_chars, page_count)
    if len(need) > cap:
        raise ParseError(
            'too_large',
            f'{len(need)} 页无文本层/文本量过低需要 OCR，超过上限 {cap} 页；'
            f'请先在源头 OCR 后再上传（避免只处理前几页却报「已入库」）',
        )
    ocr_pages, failed = [], []
    for i in need:
        try:
            blocks = ocr_fn(path, i)
        except ParseError as e:
            failed.append((i, e.payload.get('error')))
            continue
        except Exception as e:
            # PIL/fitz/子进程的意外错误不许 500 整份：按单页失败留痕（与 ParseError 同口径）
            failed.append((i, f'{type(e).__name__}: {e}'[:120]))
            continue
        if not blocks:
            # 引擎读不出字（空白页/全图无字）：按失败记，不算「已补」—— 不然 scanned 判定被骗过
            failed.append((i, 'OCR 未识别出文字'))
            continue
        out += blocks
        ocr_pages.append(i)
    # 页序：OCR 补进来的块追加在末尾，按页号排回去 —— 否则 chunk 的相邻关系（span）会错
    out.sort(key=lambda b: (b.get('page') or 0))
    if not out:
        raise ParseError('no_text_layer', '整份扫描版：文本层为空，OCR 也未补出文字')
    if scanned and not ocr_pages and \
            sum(_page_chars(b.get('text', '')) for b in out) < PDF_DOC_MIN_CHARS:
        why = failed[0][1] if failed else 'OCR 引擎不可用'
        raise ParseError('no_text_layer',
                         f'文本层共 {total} 个非空白字符，按扫描件处理；OCR 未补出文字（{why}）')
    notes = []
    if ocr_pages:
        notes.append(f'第 {", ".join(map(str, ocr_pages))} 页无文本层/文本量过低，已用 OCR 补')
    if failed:
        notes.append(f'第 {", ".join(str(i) for i, _ in failed)} 页 OCR 未成（{failed[0][1]}）')
    return notes

# 降级层**不叫 `_p_*`**：`_p_` 前缀在本文件的约定是「登记在 `CAPS` 里的入口解析器」，
# `_selftest_caps` 按这个前缀双向核对登记表与实现者。这两个是 `_p_pdf` 内部的备选层，
# 登记它们反而会让 `.pdf` 有三个「实现者」。
def _pdf_fitz(path):
    """第二级（逐页纯文本）。同样过 `_pdf_ocr_fill`：pymupdf4llm 缺席时（部分容器/宿主机），
    混合扫描件的图像页曾经在这一级**静默消失**（只出文本页、无 notes、照推 embedded）——
    与 `_p_pdf` 同一个失败族，修一处漏一处等于没修。`not out` 的 no_text_layer 由 fill 统一报。"""
    try:
        import fitz
    except Exception:
        return _pdf_pypdf(path)
    with fitz.open(path) as doc:      # with 收口：get_text 中途抛错也不泄漏句柄
        texts = [p.get_text() for p in doc]
        n = doc.page_count
    page_chars = {i: _page_chars(t) for i, t in enumerate(texts, 1)}
    out = [_blk(t.strip(), i, []) for i, t in enumerate(texts, 1) if t.strip()]
    notes = _pdf_ocr_fill(path, n, page_chars, out)
    return out, n, [], notes

def _pdf_pypdf(path):
    """BSD 许可的最后一级。三级全不可用时**复用 `_cap_pdf()` 那一句**（不再自己写一遍文案）：
    两份文案会漂移，而漂移方向恰好是「能力上报说可用、真解析报另一件事」。
    fitz 缺席才轮到这一级 —— 渲染不了页图，OCR 档只能判、不能补。"""
    try:
        from pypdf import PdfReader
    except Exception:
        raise ParseError('unsupported', _cap_pdf() or 'PDF 依赖不可用')
    r = PdfReader(path)
    try:
        texts = [(p.extract_text() or '') for p in r.pages]
    finally:
        r.close()                    # 常驻 serve 下不 close 是文件句柄泄漏（pypdf ≥ 3 支持）
    out = [_blk(t.strip(), i, []) for i, t in enumerate(texts, 1) if t.strip()]
    if not out:
        # 扫描版：与上面两级同一口径，显式失败；detail 说清这一级为什么补不了 OCR
        raise ParseError('no_text_layer',
                         '文本层为空（扫描版）；OCR 补页要 PyMuPDF 渲染页面，当前是 pypdf 兜底级')
    # 低文本量判定在这一级也要做：垃圾文本层（每页零星几个字符）按「已入库」静默过去，
    # 就是 `_pdf_ocr_fill` 文档里那个失败族 —— 这一级补不了 OCR，至少要响亮失败。
    total = sum(_page_chars(t) for t in texts)
    if _pdf_is_scanned(total, len(texts)):
        raise ParseError('no_text_layer',
                         f'文本层共 {total} 个非空白字符，按扫描件处理；'
                         'OCR 档需要 fitz 渲染页面（pip install PyMuPDF），当前不可用')
    return out, len(texts), []

_H_DOCX = re.compile(r'(?:Heading|标题)\s*(\d)')
# 伪标题（业务文档常没用 Word 样式，章节靠编号/加粗）：第X章/一、/（一）/1./1.1 等开头，
# 长句（>40 字）不像标题，排除。级别：中文序号/第X章=1，阿拉伯编号=2，整行加粗短句=3。
_PSEUDO_HEADING_MAX = 40


def _pseudo_heading_level(p, text):
    if len(text) > _PSEUDO_HEADING_MAX:
        return 0
    if re.match(r'^(?:第[一二三四五六七八九十百\d]+[章节条篇]|[一二三四五六七八九十]+[、.．]|[（(][一二三四五六七八九十]+[)）])', text):
        return 1
    if re.match(r'^\d+(?:\.\d+)*[、.．\s]', text):
        return 2
    runs = [r for r in p.runs if r.text.strip()]
    if runs and all(r.bold for r in runs):
        return 3
    return 0

def _p_docx(path):
    # 依赖不在时轮不到这里：`parse_doc` 的能力门先拒（缺依赖 = 明确 unsupported，见那里的注释）。
    # 解析器**不再各自判一次依赖**：以前每个解析器写一句自己的「缺少依赖 X」，与 PARSE_DEPS 两份真相。
    import docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    doc = docx.Document(path)
    stack, out = [], []
    for child in doc.element.body.iterchildren():      # 按正文顺序走，表格不许漏
        tag = child.tag.rsplit('}', 1)[-1]
        if tag == 'p':
            p = Paragraph(child, doc)
            if not (t := p.text.strip()):
                continue
            if m := _H_DOCX.search(getattr(p.style, 'name', '') or ''):
                _push(stack, int(m.group(1)), t)
            elif (lvl := _pseudo_heading_level(p, t)):
                # 没用 Word 样式的文档按 编号/加粗 推断章节（导图/引用定位靠这个结构）
                _push(stack, lvl, t)
            out.append(_blk(t, None, stack))
        elif tag == 'tbl':
            # 按 tblGrid 的**列位**拼，不是按「这行有几个 tc」拼：Word 允许某行晚起步/早收尾
            # （gridBefore/gridAfter），`_Row.cells` 只给实际存在的单元格 —— 不补空位的话该行整体
            # 左移一列，「标签」就跟邻行的「值」配上对（实测：开户银行 ↔ 银行账号 两行的值互换）。
            # 空单元格也**保留占位**（拼成 ` | `）：join 吞掉空串同样会毁掉列位对应关系。
            # gridSpan/vMerge 由 `_Row.cells` 自己摊平（跨列重复、竖并取上格），这里不用管。
            # getattr 取值：grid_cols_* 是 python-docx 1.2.0 才有的 API，老版本上退化成旧行为
            # 而不是 AttributeError（容器镜像的钉版见 docker/parser/Dockerfile，两边须一致）。
            rows = [' | '.join([''] * getattr(r, 'grid_cols_before', 0)
                               + [c.text.strip() for c in r.cells]
                               + [''] * getattr(r, 'grid_cols_after', 0))
                    for r in Table(child, doc).rows]
            if rows:
                out.append(_blk('\n'.join(rows), None, stack))
    return out, 0, []

def _p_pptx(path):
    from pptx import Presentation      # 依赖门在 parse_doc，同 _p_docx
    prs = Presentation(path)
    out, n = [], 0
    for i, slide in enumerate(prs.slides, 1):
        n = i
        _ts = slide.shapes.title      # shape 代理每次访问都是新对象，比身份要用 shape_id
        title = (_ts.text.strip() if _ts is not None else '') or f'第{i}页'
        title_id = _ts.shape_id if _ts is not None else None
        # 标题 shape 也有 text_frame：排除它 —— 标题进 heading_path，再进正文就被向量化两遍
        texts = [s.text_frame.text.strip() for s in slide.shapes
                 if s.has_text_frame and s.shape_id != title_id and s.text_frame.text.strip()]
        if texts:
            out.append({'text': '\n'.join(texts), 'page': i, 'heading_path': title})
    return out, n, []

def _cell(v):
    return '' if v is None else str(v).strip()   # ' 10 ' 这类带空白值不该原样进 sheets/表头

def _sheet(name, rows):
    """行列上限：超出立刻报错（不截断——截断=用户以为传成功但数据少一半；也不先吃满内存，本进程还托着 /embed）"""
    keep = []
    for r in rows:
        if not any(x != '' for x in r):
            continue
        if len(r) > MAX_COLS:
            raise ParseError('too_large', f'{name} 列数超 {MAX_COLS}（实际 {len(r)} 列）')
        keep.append(r)
        if len(keep) > MAX_ROWS + 1:
            raise ParseError('too_large', f'{name} 行数超 {MAX_ROWS}（实际至少 {len(keep)} 行）')
    if not keep:
        # 🔴 空 sheet **也要报出来**（空表头 + 空行），不能在这里丢掉。
        # Rust 侧 `tabular::plan` 的契约是「空表/无表头不建表，但把名字放进 `skipped`，
        # 不能静默 —— 用户以为整份文件都能问数，结果少了一个 sheet 是个安静的数据缺口」。
        # 它只能报**它收到过**的 sheet：在这里 return None 就等于让那条契约对「全空的 sheet」
        # 恒不成立。实测抓到：两个 sheet 的 xlsx 只回 1 个，另一个无声消失。
        # 文本通道那侧由 `tabular::sheet_blocks` 跳过它（否则会多一个只有标题的垃圾块）。
        return {'name': name, 'header': [], 'rows': []}
    return {'name': name, 'header': keep[0], 'rows': keep[1:]}

def _p_xlsx(path):
    """表格只出 sheets（单元格矩阵），markdown 文本通道由 knowledge::tabular 的 sheet_blocks 出。

    🔴 read_only 模式轻信每个 sheet 顶部的 <dimension> 声明，而 WPS/ERP 导出的 xlsx 常把它
    写小（声明 A1:A1、实际到 F 列）—— iter_rows 按声明截断，**静默丢列**（KB 审查实测形态）。
    openpyxl ≥ 3.1 用 `ws.reset_dimensions()` 让它忽略声明、按实际单元格重算边界；
    更老的版本没有它：声明列数与首行实长不符时降级非 read_only 重读（内存换正确性）。
    判据钉在 `_selftest_xlsx_dims`（篡改 dimension 声明的夹具）。"""
    import openpyxl                    # 依赖门在 parse_doc，同 _p_docx
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        out = []
        for ws in wb.worksheets:
            if hasattr(ws, 'reset_dimensions'):
                ws.reset_dimensions()  # 不信 <dimension> 声明（openpyxl ≥ 3.1）
            elif _dims_suspect(ws):
                # openpyxl < 3.1 的兜底：声明不可信 → 非 read_only 按实际单元格重读整本
                return _xlsx_eager(path)      # 本 wb 由 finally 关闭
            # `_sheet` 恒返 dict（空 sheet 返哨兵，见它的注释），旧「返 None」契约的判空已删
            out.append(_sheet(ws.title, ([_cell(c) for c in r]
                                         for r in ws.iter_rows(values_only=True))))
    finally:
        wb.close()
    return [], 0, out

def _dims_suspect(ws):
    """老 openpyxl（无 `reset_dimensions`）的错声明探测：read_only 下 `ws.max_column`
    就是 <dimension> 的声明列数，与首行实长不符即不可信。"""
    declared = ws.max_column
    if declared is None:
        return False
    first = next(ws.iter_rows(values_only=True), None)
    return first is not None and len(first) != declared

def _xlsx_eager(path):
    """非 read_only 重读整本：矩阵按实际单元格建，<dimension> 声明够不着它。"""
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=False, data_only=True)
    try:
        return [], 0, [_sheet(ws.title, ([_cell(c) for c in r]
                                         for r in ws.iter_rows(values_only=True)))
                       for ws in wb.worksheets]
    finally:
        wb.close()

def _p_csv(path):
    text = _read_text(path)
    try:
        dialect = csv.Sniffer().sniff(text[:4096], delimiters=',;\t|')
    except csv.Error:
        dialect = csv.excel
    rows = ([_cell(c) for c in r] for r in csv.reader(io.StringIO(text), dialect))
    return [], 0, [_sheet(os.path.splitext(os.path.basename(path))[0], rows)]

OCR_LANG = os.environ.get('DMS_OCR_LANG', 'chi_sim+eng')

def _p_image(path):
    """图片 OCR → 一个块（图片没有章节结构，`heading_path` 用文件名，引用时看得出来源）。

    🔴 **识别不出文字必须报错**，不许 `return [], 0, []`。静默返空的后果不是「少一份文档」：
    Rust `ingest::run` 拿到 0 块会走 `chunks.is_empty()` 那条 BadInput，但只要同一份文件
    还带出过 sheets 就会照旧推到 `embedded` —— 界面显示「已入库」、chunk_count=0、
    问什么都答不出来。那是本仓反复抓的静默失败族，不许在新增格式时再造一遍。
    `no_text_layer` 与扫描版 PDF 同一个码：Rust 侧已把它映成「需先 OCR」的确定性失败。

    🔴 **OCR 引擎：千问 flash 优先，tesseract 降级**（业主 2026-07-31 裁决「图片识别用千问」）。
    实测（`_silent/` 三个扫描件，全部只读实测）：
      - `multiframe.tif` 两帧 → 千问 689/712ms 全对（`TIFFPAGE2-7788`）；tesseract 只认第一帧
      - `scanned.pdf` → 千问 896ms 读出 `SCANONLY-3344`；tesseract 产空文本
      - `mixed.pdf` 文本层+图像页 → 千问逐页全对（`TEXTPAGE1-5566` / `PDFOCR2-9911`）
    千问 flash 是 vision 模型（自己吃图，988ms 级，3/3 全对，比 tesseract 准且**逐帧不丢**）。
    零新依赖（`urllib.request` 发 HTTP，不引 openai 包 —— 宿主机 SAC 会拦新编译扩展）。
    配置：`DMS_QWEN_OCR_KEY`（或 `QWEN_KEY`，只读环境变量、不回退 settings 的 llm_api_key）、
    `DMS_QWEN_OCR_MODEL`（默认 qwen3.7-flash）、
    `DMS_QWEN_OCR_BASE`（默认 dashscope compatible-mode）。千问不可用时回落 tesseract。
    """
    from PIL import Image, ImageSequence
    name = os.path.basename(path)
    blocks, frames = [], 0
    try:
        with Image.open(path) as im:
            # 帧数护栏（与 PDF 侧 OCR_PAGE_CAP 同口径）：每帧一次千问 HTTP（60s 超时），
            # N 帧最坏远超 Rust 120s 解析超时 —— 超了响亮 too_large，不「OCR 前 N 帧报已入库」
            n_frames = getattr(im, 'n_frames', 1)
            if n_frames > OCR_PAGE_CAP:
                raise ParseError('too_large',
                                 f'图片 {n_frames} 帧超上限 {OCR_PAGE_CAP}（每帧一次 OCR 请求）')
            for i, frame in enumerate(ImageSequence.Iterator(im), 1):
                frames = i
                # 优先千问；不可用/失败回落 tesseract（两路同一形状：一帧一块）
                t = _ocr_qwen_frame(frame) or _ocr_tesseract_frame(frame)
                if t:
                    # 多帧各带帧号：共用文件名做 heading_path 会被 chunk 分组合并出跨帧重复
                    hp = name if n_frames == 1 else f'{name}#f{i}'
                    blocks.append({'text': t, 'page': i, 'heading_path': hp})
    except ParseError:
        raise            # ParseError 原样上抛：再包一层 detail 双重套娃（OCR 失败（…OCR 失败（…）））
    except Exception as e:
        raise ParseError('unsupported', f'OCR 失败（{e}）')
    if not blocks:
        raise ParseError('no_text_layer', f'图片 OCR 未识别出文字（{frames} 帧全空）')
    notes = [] if len(blocks) == frames else [
        f'第 {i} 帧未识别出文字' for i in range(1, frames + 1)
        if not any(b['page'] == i for b in blocks)
    ]
    return blocks, frames, [], notes


def _ocr_tesseract_frame(frame):
    """tesseract 那一帧。`ParseError` 留给 `_p_image` 统一处理（语言包缺失要报出来）。"""
    import pytesseract
    if exe := _exe('DMS_TESSERACT', 'tesseract'):
        pytesseract.pytesseract.tesseract_cmd = exe
    try:
        return pytesseract.image_to_string(frame, lang=OCR_LANG).strip()
    except Exception as e:
        raise ParseError('unsupported', f'tesseract OCR 失败（lang={OCR_LANG}）：{e}')


def _ocr_qwen_frame(frame):
    """千问 flash 那一帧。**任何失败都返 None**（由 `_p_image` 回落 tesseract），不抛 —
    一张图的网络抖动不该让整份扫描件失败（回落的方向是「降级但入库」，不是「整份拒收」）。
    返回空串也算 None：模型读不出字时静默回落，让 tesseract 再试一次。"""
    key = os.environ.get('DMS_QWEN_OCR_KEY') or os.environ.get('QWEN_KEY', '')
    if not key:
        return None
    import base64
    b = io.BytesIO()                 # io 顶层已有，不重复 import
    frame.convert('RGB').save(b, 'PNG')
    b64 = base64.b64encode(b.getvalue()).decode()
    body = {
        'model': os.environ.get('DMS_QWEN_OCR_MODEL', 'qwen3.7-flash'),
        'temperature': 0.1,
        'enable_thinking': False,
        'messages': [{'role': 'user', 'content': [
            {'type': 'text', 'text':
                '读出这张图里的全部文字。如果是一则通知/制度/标准，'
                '把标题、要点、数字、金额、日期逐条列出。只输出读到的内容，不要解释。'},
            {'type': 'image_url', 'image_url': {'url': f'data:image/png;base64,{b64}'}},
        ]}],
    }
    base = os.environ.get(
        'DMS_QWEN_OCR_BASE', 'https://dashscope.aliyuncs.com/compatible-mode/v1')
    try:
        req = urllib.request.Request(
            f'{base}/chat/completions', data=json.dumps(body).encode(),
            headers={'Authorization': f'Bearer {key}', 'Content-Type': 'application/json'})
        with urllib.request.urlopen(req, timeout=60) as r:
            v = json.load(r)
        return (v['choices'][0]['message'].get('content') or '').strip() or None
    except Exception:
        return None

IMG_EXTS = ('.png', '.jpg', '.jpeg', '.bmp', '.tif', '.tiff', '.webp', '.gif')
# ── 扫描件检测与 OCR 档（Y1）：文本层「低文本量」即判扫描件 → 页面渲染成图 → vision OCR ──
# 阈值是**可调常量**（环境变量覆盖），判据由 `_selftest_pdf_scan` 钉住 —— 改阈值先改钉。
def _env_int(name, default):
    """阈值环境变量的整数解析：写错值要指出是哪个变量、值是什么 ——
    裸 ValueError traceback 会让服务起不来且不知所云。"""
    v = os.environ.get(name)
    if v is None:
        return default
    try:
        return int(v)
    except ValueError:
        raise SystemExit(f'环境变量 {name}={v!r} 不是整数，embed 服务无法启动')

# 单页非空白字符 < 50 → 该页送 OCR（零文本页是它的特例：原「没有块的页」判据被它覆盖）。
PDF_PAGE_MIN_CHARS = _env_int('DMS_PDF_PAGE_MIN_CHARS', 50)
# 整份判扫描件：全文 < 200，或页均 < 单页阈值（垃圾文本层：每页零星几个字符，合计不少、内容没有）。
PDF_DOC_MIN_CHARS = _env_int('DMS_PDF_DOC_MIN_CHARS', 200)
# 渲染 dpi（`_pdf_page_ocr`）：200 是实测折中 —— 150 下 12pt 中文 tesseract 常丢字，300 单页多花 ~1.5 倍。
OCR_DPI = _env_int('DMS_OCR_DPI', 200)
# PDF 逐页补 OCR 的页数上限（`_pdf_ocr_fill`）。可用 DMS_OCR_PAGE_CAP 覆盖。
# 成本口径：OCR 档每页 ≈ fitz 渲染 0.1s（200dpi）+ 识别 0.2~1s（千问 vision ~1s/页，tesseract 0.2~1s/页）。
# 30 页 × ~1s ≈ 半分钟，仍低于 Rust 侧 120s 解析超时（connector/src/doc.rs PARSE_TIMEOUT_SECS）；
# 超 cap 不「OCR 前 N 页然后报已入库」—— too_large 响亮失败。
OCR_PAGE_CAP = _env_int('DMS_OCR_PAGE_CAP', 30)
LEGACY_TARGET = {'.doc': '.docx', '.xls': '.xlsx', '.ppt': '.pptx'}
SOFFICE_TIMEOUT = 120      # 每次转换都用全新临时 profile（`_p_legacy` 为并发隔离），建 profile 的成本每次都付；超时要响亮而不是挂死请求

def _p_legacy(path):
    """旧二进制 Office（OLE/CFB 的 .doc/.xls/.ppt）：**不自己解**，交 LibreOffice headless
    转成对应的新格式，再走已有的 `_p_docx/_p_xlsx/_p_pptx`（一条解析通道，不是两套代码）。

    `-env:UserInstallation` 给每次调用一个独立 profile：共享默认 profile 时，第二个并发
    soffice **直接退出且什么都不产出，返回码仍是 0** —— 那正好是「转换成功但没有文件」
    这种静默形态。本函数据此只信「产物文件在不在」，不信返回码。
    soffice 缺席时轮不到这里：`parse_doc` 的能力门先拒（`_cap_legacy` 同时要求目标格式的解析器，
    少查那一半就会变成「转出来了却解析不了」）。
    """
    ext = os.path.splitext(path)[1].lower()
    tgt = LEGACY_TARGET[ext]
    with tempfile.TemporaryDirectory() as d:
        prof = 'file:///' + os.path.join(d, 'profile').replace('\\', '/').lstrip('/')
        try:
            r = subprocess.run([_soffice(), f'-env:UserInstallation={prof}', '--headless',
                                '--convert-to', tgt[1:], '--outdir', d, path],
                               capture_output=True, timeout=SOFFICE_TIMEOUT)
        except subprocess.TimeoutExpired:
            raise ParseError('unsupported', f'soffice 转换 {ext} 超过 {SOFFICE_TIMEOUT}s')
        out = os.path.join(d, os.path.splitext(os.path.basename(path))[0] + tgt)
        if not os.path.isfile(out):
            why = (r.stderr or r.stdout or b'').decode('utf-8', 'replace').strip()[:200]
            raise ParseError('unsupported', f'soffice 未能把 {ext} 转成 {tgt}（rc={r.returncode}）{why}')
        return CAPS[tgt][0](out)

# ── 依赖探测：每种格式「装了没」的**唯一实现**。返回 '' = 可用，否则返回给运维看的原因 ──
# 🔴 原先是两份真相：`PARSE_DEPS` 一份、每个解析器自己那句「缺少依赖 python-docx」一份。
# 两份文案会漂移，而漂移方向恰好是**上报说可用、真解析报另一件事**（`find_spec` 那次就是这么发生的）。
# 现在解析器不判依赖，`parse_doc` 在入口过这道门 —— 能力上报与拒绝理由必然同一句话。

def _mod(mod, pip):
    """单模块依赖。`mod` 是 **import 名**（python-docx 的 import 名是 docx），`pip` 是装包名。"""
    return lambda: '' if _have(mod) else f'{mod} 不可用（pip install {pip}）：{_why(mod)}'

def _cap_ok():
    return ''      # md/txt/csv 纯标准库，恒可用

def _cap_pdf():
    """三级任一即可（见 `_p_pdf`）。装到哪一级决定有没有章节层级，由 `_pdf_tier` 照实上报。"""
    if _have('pymupdf4llm') or _have('fitz') or _have('pypdf'):
        return ''
    return ('PDF 三级依赖全不可用（pip install pymupdf4llm —— AGPL-3.0，业主已裁决可用，'
            '带章节层级 heading_path；或 pypdf —— BSD-3，只出逐页纯文本）：'
            + '；'.join(f'{m}（{_why(m)}）' for m in ('pymupdf4llm', 'fitz', 'pypdf')))

def _pdf_text_engine():
    """text 档三级引擎里当前会用的那一级（`_pdf_tier` 的人话与 `tiers` 的机读值同源，不许漂）。"""
    return ('pymupdf4llm' if _have('pymupdf4llm')
            else 'fitz' if _have('fitz') else 'pypdf')

def _ocr_engine():
    """OCR 档实际会用的引擎：'qwen' / 'tesseract' / None。
    与 `_cap_ocr` 同一判定、只是报名字 —— 两处判定不许漂（上报可用、真解析不可用 = 两套口径）。"""
    if (os.environ.get('DMS_QWEN_OCR_KEY') or os.environ.get('QWEN_KEY')) and _have('PIL.Image'):
        return 'qwen'
    if _have('PIL.Image') and _have('pytesseract') and _exe('DMS_TESSERACT', 'tesseract'):
        return 'tesseract'
    return None

def _pdf_tier():
    """text 档 + ocr 档的自报人话（Y1：扫描件高阶档）。ocr 档要**两半**：渲染（fitz 把页
    渲成图）+ 识别引擎（千问 vision / tesseract）—— 缺一半就是不可用，自报必须说清缺哪半。"""
    text = {'pymupdf4llm': 'pymupdf4llm：带章节层级 heading_path',
            'fitz': 'fitz 兜底：只有逐页纯文本，heading_path 为空',
            'pypdf': 'pypdf 兜底（BSD-3）：只有逐页纯文本，heading_path 为空'}[_pdf_text_engine()]
    eng = _ocr_engine()
    if not _have('fitz'):
        ocr = '不可用（fitz 缺席，扫描件页面渲不成图；pip install PyMuPDF）'
    elif eng == 'qwen':
        ocr = '千问 vision flash 优先，tesseract 降级'
    elif eng == 'tesseract':
        ocr = 'tesseract（本地）'
    else:
        ocr = '不可用（OCR 引擎缺席：配 DMS_QWEN_OCR_KEY 走千问，或装 tesseract，详见 .png 能力行）'
    return f'{text}；扫描件 OCR 档：{ocr}'

def _cap_ocr():
    """OCR 引擎：**千问 flash 或 tesseract 有一样就能用**（`_p_image` 优先千问、降级 tesseract）。
    千问要 `DMS_QWEN_OCR_KEY`（或 `QWEN_KEY`）+ Pillow（帧转 png 用）。
    tesseract 要三样：Pillow、pytesseract（只是包装器）、tesseract **二进制**。
    探 `PIL.Image` 而不是 `PIL`：编译扩展 `_imaging` 在子模块里，正是本机 SAC 拦的那一类。"""
    qwen_ok = bool(os.environ.get('DMS_QWEN_OCR_KEY') or os.environ.get('QWEN_KEY'))
    if qwen_ok and _have('PIL.Image'):
        return ''       # 千问路通，tesseract 那套不是必需
    miss = [(m, pip) for m, pip in (('PIL.Image', 'pillow'), ('pytesseract', 'pytesseract'))
            if not _have(m)]
    if miss:
        # 按实际缺的那半给安装建议：千问路根本不需要 pytesseract，一句「全装」误导运维多装
        pips = ' '.join(dict.fromkeys(pip for _, pip in miss))
        return f'OCR 依赖不可用（pip install {pips}）：' + '；'.join(f'{m}（{_why(m)}）' for m, _ in miss)
    if not _exe('DMS_TESSERACT', 'tesseract'):
        return ('找不到 tesseract 可执行文件（pytesseract 只是包装器，二进制要单独装）：'
                'apt-get install tesseract-ocr tesseract-ocr-chi-sim，或 DMS_TESSERACT 指全路径；'
                '或配 DMS_QWEN_OCR_KEY 走千问 flash（业主裁决，图片识别用千问）')
    return ''

def _cap_legacy(target):
    """旧二进制 Office 要**两半**：soffice 转格式 + 目标格式的解析器（.doc→.docx 仍要 python-docx）。"""
    return lambda: ('找不到 soffice/libreoffice 可执行文件（旧二进制 Office 靠 LibreOffice headless '
                    '转成新格式）：apt-get install libreoffice-writer/-calc/-impress，'
                    '或 DMS_SOFFICE 指全路径') if not _soffice() else CAPS[target][2]()

# 🔴 **单一事实源**：扩展名 → (解析器, 能力名, 依赖探测)。
# 以前是两张表（`PARSERS` 管扩展名、`PARSE_DEPS` 管依赖），于是「登记而不消费」查不出来：
# 声明支持某扩展名却没有实现者、或写了解析器却没登记（/parse 永远走不到），两种都静默。
# 现在 `parse_doc` / `parse_ok` / `parse_caps` 全从这一张表派生，`_selftest_caps` 双向钉住。
CAPS = {
    '.pdf': (_p_pdf, 'pdf', _cap_pdf),
    '.docx': (_p_docx, 'docx', _mod('docx', 'python-docx')),
    '.pptx': (_p_pptx, 'pptx', _mod('pptx', 'python-pptx')),
    '.xlsx': (_p_xlsx, 'xlsx', _mod('openpyxl', 'openpyxl')),
    '.xlsm': (_p_xlsx, 'xlsx', _mod('openpyxl', 'openpyxl')),
    '.csv': (_p_csv, 'text', _cap_ok),
    '.md': (_p_text, 'text', _cap_ok),
    '.markdown': (_p_text, 'text', _cap_ok),
    '.txt': (_p_text, 'text', _cap_ok),
    # json/log/html：纯文本族，零第三方依赖（html 去标签走标准库 HTMLParser）
    '.json': (_p_json, 'text', _cap_ok),
    '.log': (_p_text, 'text', _cap_ok),
    '.html': (_p_html, 'text', _cap_ok),
    # .htm 刻意不登记：Rust 侧 EXTS 白名单（knowledge/src/ingest.rs）与本表是**集合相等**判据，
    # 单侧加 .htm 会红那边的 `exts_cover_the_doc_service_capabilities`；要支持需两侧同加（本轮只动本文件，暂缓）。
    # 旧二进制 Office：各自一个能力名（三个 LibreOffice 组件是分开装的，合成一个键会
    # 让「只装了 writer」上报成整族可用）
    '.doc': (_p_legacy, 'doc', _cap_legacy('.docx')),
    '.xls': (_p_legacy, 'xls', _cap_legacy('.xlsx')),
    '.ppt': (_p_legacy, 'ppt', _cap_legacy('.pptx')),
    **{e: (_p_image, 'image', _cap_ocr) for e in IMG_EXTS},
}
MIME_EXT = {
    'application/pdf': '.pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
    'application/msword': '.doc',
    'application/vnd.ms-excel': '.xls',
    'application/vnd.ms-powerpoint': '.ppt',
    'image/png': '.png', 'image/jpeg': '.jpg', 'image/bmp': '.bmp',
    'image/tiff': '.tif', 'image/webp': '.webp', 'image/gif': '.gif',
    'text/csv': '.csv', 'text/markdown': '.md', 'text/plain': '.txt',
    'application/json': '.json', 'text/html': '.html',
}

def parse_doc(path, mime=''):
    ext = os.path.splitext(path)[1].lower()
    if ext not in CAPS:
        # mime 大小写不敏感（RFC 允许 `Application/PDF` 这种写法），查表前统一 lower
        ext = MIME_EXT.get((mime or '').split(';')[0].strip().lower(), ext)
    cap = CAPS.get(ext)
    if cap is None:
        raise ParseError('unsupported', f'不支持的格式：{ext or mime or path}')
    if not os.path.isfile(path):
        raise ParseError('not_found', path, 404)
    # 依赖门**只有这一处**：缺依赖 = 明确的 `unsupported` + 一句原因（Rust 侧映成 BadInput，
    # 文档落 failed 且 error 里写清缺什么）。绝不许让解析器悄悄返回空 blocks —— 那条路的终点是
    # 「status=embedded、chunk_count=0、界面显示已入库、问什么都答不出来」。
    if why := cap[2]():
        raise ParseError('unsupported', f'{ext} 暂不可用：{why}')
    got = cap[0](path)
    # 解析器可以返 3 元或 4 元。**第三位恒是 `sheets`（表格类文件的工作表）**，
    # 第四位是可选的 `notes`（人话说明，例：「第 2 页无文本层，已用 OCR 补」）。
    #
    # 🔴 为什么必须分开：`sheets` 在 Rust 侧是 `Vec<Sheet>`（`connector/src/doc.rs::ParsedDoc`），
    # 往里塞字符串会让**整份响应反序列化失败**。我第一版把 PDF 的页注塞进第三位，
    # 混合 PDF 在 Rust 路径上就会直接解析失败 —— 自己实测响应体时才看出来
    # （`sheets = "第 2 页无文本层，已用 OCR 补"`）。
    # `notes` 是新键；Rust 侧没有 `deny_unknown_fields`，多这个键不影响老逻辑，
    # 将来要在界面上显示它只需加一个字段。
    blocks, page_count, sheets = got[0], got[1], got[2]
    out = {'blocks': blocks, 'page_count': page_count, 'sheets': sheets}
    if len(got) > 3 and got[3]:
        out['notes'] = got[3]
    return out

_HAVE_CACHE = {}
_HAVE_ERR = {}
_EXE_CACHE = {}

def _exe(env, *names):
    """外部**可执行文件**探测：soffice / tesseract 不是 python 模块，`_have` 探不到它们。
    `env` 优先（容器里可能装在非 PATH 位置）。缓存同 `_have`（`/health` 每次都调）——
    代价是改了 PATH/环境变量要重启服务，与依赖装卸同一个口径。"""
    key = (env, names)         # 键带 names：同名 env 不同 names 的两个调用点不许互相污染缓存
    if key not in _EXE_CACHE:
        _EXE_CACHE[key] = os.environ.get(env) or next(
            (p for n in names if (p := shutil.which(n))), '')
    return _EXE_CACHE[key]

def _soffice():
    return _exe('DMS_SOFFICE', 'soffice', 'libreoffice')

def _have(mod):
    """**真 import 一次**，不是 `find_spec`。

    🔴 `find_spec` 只查「包在不在」，查不出「import 会不会炸」。实测踩到：
    `python-docx` / `python-pptx` 装好了（spec 在），但它们依赖的 `lxml` 编译扩展被本机
    Smart App Control 拦掉（`DLL load failed while importing etree`）→ `find_spec` 说可用、
    真解析当场抛。那正是本函数文档里那句「不许假装可用」要防的事，而它自己犯了。
    结果缓存：`/health` 每次都调，重复 import 走 `sys.modules` 已经很快，但异常路径不缓存
    会每次重跑一遍失败的 DLL 加载。
    """
    if mod not in _HAVE_CACHE:
        try:
            importlib.import_module(mod)
            _HAVE_CACHE[mod] = True
        except Exception as e:      # ImportError / OSError(DLL) / 任何加载期异常都算不可用
            _HAVE_CACHE[mod] = False
            _HAVE_ERR[mod] = f'{type(e).__name__}: {e}'[:200]
    return _HAVE_CACHE[mod]

def _why(mod):
    """`_have` 失败的**原文**。必须进能力上报：不然「没装」与「装了但加载不了」长得一样，
    而处置完全相反。本机 `.venv` 实测两种都有：
      docx/pptx → `ImportError: DLL load failed while importing etree: 应用程序控制策略已阻止此文件。`
                  （装了，是 lxml 的编译 DLL 被 Smart App Control 拦）
      pymupdf4llm → `ModuleNotFoundError: No module named 'pymupdf4llm'`（真没装）
    只报一句「pip install python-docx」会让运维对着第一种白装一遍，再回来问为什么还是红的。"""
    _have(mod)
    return _HAVE_ERR.get(mod, '')

def parse_ok():
    """能力名 → 是否真的可用（粗粒度，历史形状）。`tools/parse_probe.py` 按 `pdf`/`xlsx`
    两个键决定跳不跳夹具 —— 键改名它会**静默全跳过并 exit 0**（判据恒绿），
    所以 `_selftest_caps` 有一条断言钉着这两个键。逐扩展名的原因看 `parse_caps()`。
    同名能力取**与**（不是任取一个）：一个能力名下的扩展名探测可以不同（.doc/.xls/.ppt 就是），
    落到 `dict` 推导里会被后写的覆盖，那正是「上报可用、其实只装了一半」。"""
    ok = {}
    for _, fam, probe in CAPS.values():
        ok[fam] = ok.get(fam, True) and not probe()
    return ok

def parse_caps():
    """**逐扩展名**的能力上报：`{'.docx': {'ok': bool, 'why': '一句人话'}}`。
    判定只看 `ok`；`why` 给人看 —— 不可用时说缺什么怎么装，可用时可能说降级（见下面 pdf 那句）。"""
    caps = {ext: {'ok': not (w := probe()), 'why': w} for ext, (_, _, probe) in CAPS.items()}
    # pdf 三级依赖里装到哪一级决定有没有章节层级，而那是升到 pymupdf4llm 的**全部理由**。
    # 上报里不说，运维就无从知道自己拿到的是「带章节」还是「纯文本」那种质量。
    if caps['.pdf']['ok']:
        caps['.pdf']['why'] = _pdf_tier()
        # 机读两档（Y1）：text = 三级文本引擎；ocr = 扫描件 OCR 引擎（fitz 缺席时恒 None，
        # 渲染那半先断了）。自报与真解析同源 —— `_pdf_ocr_fill` 用的就是这两样。
        caps['.pdf']['tiers'] = {'text': _pdf_text_engine(),
                                 'ocr': _ocr_engine() if _have('fitz') else None}
    return caps

# ============ 文档服务：分块（K1）============
CHARS_PER_TOKEN = 1.6      # 中文口径：1.6 字符/token（不是英文那套 4 字符/token，照搬会切出两倍大的块）
TARGET_TOKENS, OVERLAP = 400, 60
MAX_TOKENS = 480           # 硬上限：bge-small-zh-v1.5 max_seq_len=512，超了 fastembed 静默截断（症状是检索时好时坏）
_SENT = re.compile(r'(?<=[。！？；!?;\n])')

def est_tokens(text):
    return math.ceil(len(text) / CHARS_PER_TOKEN)

def _split_long(text, cap):
    """超长块先按句末标点切，仍超长则硬切"""
    parts, buf = [], ''
    for piece in _SENT.split(text):
        if buf and len(buf) + len(piece) > cap:
            parts.append(buf)
            buf = ''
        while len(piece) > cap:
            parts.append(piece[:cap])
            piece = piece[cap:]
        buf += piece
    if buf.strip():
        parts.append(buf)
    return parts

def _emit(chunks, text, hp, page):
    if not (text := text.strip()):
        return
    n = est_tokens(text)
    # 不用 assert：python -O 会把 assert 整条剥掉，超窗块静默进库被 fastembed 截断
    if n > MAX_TOKENS:
        raise ValueError(f'块 {n} token 超上限 {MAX_TOKENS}（bge 512 窗口会静默截断）')
    chunks.append({'text': text, 'heading_path': hp, 'page': page, 'tokens': n})

def _fill(chunks, blocks, hp, target_chars, overlap_chars, cap):
    """同一 heading_path 内：按目标长度合并/切分，块间重叠 overlap_chars。

    🔴 `page` 由 buf **吸收过的页集合**定，跨页就置 None。
    原实现是 `if page is None: page = b.get('page')` —— 只取**第一个**贡献块的页，
    而 buf 之后还会吸收后面几页的块（PDF 的 heading stack 跨页延续，见 `_p_pdf`：
    一个 heading_path 的块常横跨 2-3 页）⇒ 一个跨 3 页的块标着「第 1 页」，
    用户翻到第 1 页找不到引用的那句话，还会以为引用是编的。
    `answer.rs::source_of` 与 `KbAnswer.vue::loc` 都在 page 为 None 时不显示页码：
    **「不知道」比「说错」好**。判据见 `_selftest_pages`。
    重叠尾巴不算跨页（它只是上一块末尾 overlap_chars 个字符的上下文，不是本块的实质内容）；
    `pages.add` 必须在 flush **之后**，否则触发 flush 的那一页会污染刚发出去的块。
    """
    def one_page(pages):
        real = {p for p in pages if p is not None}   # 无页码的块（md）不算一页
        return real.pop() if len(real) == 1 else None

    buf, pages, fresh = '', set(), False
    for b in blocks:
        for unit in _split_long((b.get('text') or '').strip(), cap):
            if buf and len(buf) + 1 + len(unit) > target_chars:
                # ponytail: 重叠取字符尾巴（可能从句中开始），检索质量不满意再改成按句边界回退
                _emit(chunks, buf, hp, one_page(pages))
                buf, pages, fresh = buf[len(buf) - overlap_chars:] if overlap_chars else '', set(), False
            pages.add(b.get('page'))
            buf = (buf + '\n' + unit) if buf else unit
            fresh = True
    if fresh:
        _emit(chunks, buf, hp, one_page(pages))

def _selftest_pages():
    """`_fill` 的页码判据，三条都反向验证过（2026-07-30，实际打坏再跑 `selftest`）：
    - `one_page` 换回「取第一个非空页」（`min(real)`）→ `across` 那条红：跨 3 页的块 page==1（原 bug）
    - `one_page` 一律返回 None → `same` 那条红：单页块也丢了页码（别把修法做成「一律不显示」）
    - `pages.add` 挪到 flush **之前** → `split` 那条红：两块都 None（第 2 页串味给刚发出去的第 1 块）"""
    same = [{'text': '甲' * 20, 'heading_path': 'H', 'page': 1},
            {'text': '乙' * 20, 'heading_path': 'H', 'page': 1}]
    across = [{'text': '甲' * 20, 'heading_path': 'H', 'page': 1},
              {'text': '乙' * 20, 'heading_path': 'H', 'page': 2},
              {'text': '丙' * 20, 'heading_path': 'H', 'page': 3}]
    a, b = chunk_blocks(same), chunk_blocks(across)
    assert len(a) == 1 and len(b) == 1, (a, b)        # 都合成一块，唯一差别就是页
    assert a[0]['page'] == 1, a
    assert b[0]['page'] is None, b                    # 跨 3 页 → 宁可不显示
    # 无页码的块（md）不算一页：单页 + 无页码仍按那唯一的真页算（保持原行为）
    mix = chunk_blocks([{'text': '甲' * 20, 'heading_path': 'H'},
                        {'text': '乙' * 20, 'heading_path': 'H', 'page': 7}])
    assert len(mix) == 1 and mix[0]['page'] == 7, mix
    # 切成多块时各记自己的页（别把「跨页置 None」做成「一律 None」，也别让 flush 那一页串味）
    split = chunk_blocks([{'text': '甲' * 400, 'heading_path': 'H', 'page': 1},
                          {'text': '乙' * 400, 'heading_path': 'H', 'page': 2}], overlap=0)
    assert [c['page'] for c in split] == [1, 2], split

# `| --- | :-: |` 形分隔行：单元格只许是 `-` 与 `:`（带数字/文字的行不是分隔行）
_MD_SEP = re.compile(r'^\s*\|?(?:\s*:?-+:?\s*\|)+\s*$')

def _table_split(text):
    """markdown 表格块的「表头切点」：返回 (每块都要重复的前缀, 数据行列表)；不是表格块返 None。

    只认**块首 6 行内**的表格（sheet 渲染块恒在块首：`# 名` + 空行 + 表头行 + 分隔行）——
    块中间的表格仍走通用分块：那种前缀里裹着散文，「每块重复前缀」会把散文复制 N 份，更糟。"""
    lines = text.split('\n')
    for i, l in enumerate(lines[:6]):
        if i > 0 and _MD_SEP.match(l) and lines[i - 1].lstrip().startswith('|'):
            return '\n'.join(lines[:i + 1]), [x for x in lines[i + 1:] if x.strip()]
    return None

def _fill_table(chunks, prefix, rows, hp, page, tc):
    """表格块的**行感知**装箱（KB 审查实测：通用分块从表中间切开，后半表是没表头的裸数字）。
    按 `\n` 行装箱到目标长度，**每块重复「# 标题 + 表头 + 分隔行」前缀**（markdown 表跨块的
    标准做法）；块间不做字符重叠 —— 表格的重叠尾巴是半行裸数据，表头重复才是上下文。"""
    budget = tc - len(prefix) - 1                            # 每块数据行的字符预算
    # budget ≤ int(480×1.6)−len(prefix)−1：预算本身就破不了 512 窗口（原 hard_cap 恒 ≥ budget，死防御已删）
    cur, cur_len = [], 0
    for row in rows:
        need = len(row) + 1
        if cur and cur_len + need > budget:
            _emit(chunks, prefix + '\n' + '\n'.join(cur), hp, page)
            cur, cur_len = [], 0
        if need > budget:
            # 单行就超目标（单元格超长的病态行）：带表头硬切，仍不破 MAX_TOKENS 窗口（见上面 budget 注释）
            for piece in _split_long(row, max(1, budget)):
                _emit(chunks, prefix + '\n' + piece, hp, page)
            continue
        cur.append(row)
        cur_len += need
    if cur:
        _emit(chunks, prefix + '\n' + '\n'.join(cur), hp, page)
    elif not rows and prefix.strip():
        # 纯表头表（0 数据行的表单模板/空表）：表头本身就是内容（字段清单就是用户会问的），
        # 这里不发块 = 整份文档 0 块、入库报「没有可索引的文本」（实测《巡店记录表单-模板.xlsx》）。
        _emit(chunks, prefix, hp, page)

def chunk_blocks(blocks, target_tokens=TARGET_TOKENS, overlap=OVERLAP):
    overlap = max(0, min(overlap, MAX_TOKENS // 4))
    target = max(1, min(target_tokens, MAX_TOKENS - overlap))
    tc, oc = int(target * CHARS_PER_TOKEN), int(overlap * CHARS_PER_TOKEN)
    cap = max(1, tc - oc - 1)   # 单元留出重叠余量：短标题块才能与正文合并，且 重叠+单元 不破 MAX_TOKENS
    chunks = []
    for hp, group in itertools.groupby(blocks, key=lambda b: b.get('heading_path') or ''):
        prose = []
        for b in group:
            t = (b.get('text') or '').strip()
            split = _table_split(t) if t else None
            if split is None or tc - len(split[0]) - 1 < 50:
                # 非表格块照旧；表头本身就吃掉大半预算的宽表（如 200 列）也回通用路径。
                # 50 的依据：数据行连 50 字符预算都分不到时，行感知装箱每块只放得下一行，
                # 不如回通用路径（与 PDF_PAGE_MIN_CHARS 同个「一点残量不算内容」的量级）
                prose.append(b)
                continue
            _fill(chunks, prose, hp, tc, oc, cap)   # 先收掉表格前面的散文段
            prose = []
            _fill_table(chunks, split[0], split[1], hp, b.get('page'), tc)
        _fill(chunks, prose, hp, tc, oc, cap)
    return chunks

def _vlit(v):
    """pgvector 字面量（与 Rust 侧 embed::to_pgvector 同一形状）"""
    return '[' + ','.join(f'{x:.6f}' for x in v) + ']'

def _revec(cur, what, sel, upd, sel_params=(), upd_tail=(), is_query=False):
    """两列 (key, 向量化文本) 的 SELECT → embed → 按 key 回写 embedding，返回条数。
    🔴 sel/upd 的 ds 限定由调用方拼死：对应 Rust 侧 `meta::DS_PRED`（单一事实源，
       `crates/server/src/meta.rs`）。python 侧不受那条漂移守卫保护 —— 漏一处就是跨源
       乱改 embedding，而 embedding 正是表召回与选源的依据。
    空结果直接返回：连模型都不加载（离线重跑常态是「没有待处理行」）。"""
    cur.execute(sel, sel_params)
    rows = cur.fetchall()
    if not rows:
        return 0
    print(f'计算 {len(rows)} {what} embedding …', flush=True)
    vecs = embed([(r[1] or '')[:1000] for r in rows], is_query)
    if len(vecs) != len(rows):
        # 条数不符不许 zip 静默截断：少返一行会把向量错位回写到别的行（Rust embed.rs 同一纪律）
        raise RuntimeError(f'{what} embedding 返回 {len(vecs)} 条，与 {len(rows)} 行不符')
    for r, v in zip(rows, vecs):
        cur.execute(upd, (_vlit(v), r[0]) + upd_tail)
    return len(rows)

def build(ds='dms'):
    import psycopg2
    # connect_timeout 与 revec 同口径：PG 不在时 5s 响亮失败，不挂死在 TCP 上
    pg = psycopg2.connect(connect_timeout=5, **pg_conf()); pg.autocommit = True; cur = pg.cursor()
    try:
        cur.execute("CREATE EXTENSION IF NOT EXISTS vector")
        cur.execute(f"ALTER TABLE meta.table_doc ADD COLUMN IF NOT EXISTS embedding vector({DIM})")
        # table_doc 刻意不带 embedding IS NULL 过滤（exemplar/element 的 SELECT 都带）：
        # search_doc 变更的置 NULL 失效钩子不在本服务视野内，全量重算是防「文本改了向量没改」
        # 的保守口径；代价是每次 build 重算全部表向量，接受它。
        n_tbl = _revec(
            cur, '表',
            "SELECT table_name, coalesce(nullif(search_doc, ''), table_name) FROM meta.table_doc"
            " WHERE ds_id = %s",                      # ds 限定 = Rust 侧 meta::DS_PRED
            "UPDATE meta.table_doc SET embedding = %s WHERE table_name = %s AND ds_id = %s",
            (ds,), (ds,))
        cur.execute("DROP INDEX IF EXISTS meta.idx_doc_hnsw")
        # 上面刚 DROP 过，这里裸 CREATE INDEX：再写 IF NOT EXISTS 读起来像不确定 drop 过没有
        cur.execute("CREATE INDEX idx_doc_hnsw ON meta.table_doc"
                    " USING hnsw (embedding vector_cosine_ops)")
        # 语料问句向量（供语义缓存召回；问句侧 query embedding，与 Rust 回写一致）
        n_ex = _revec(
            cur, '条语料问句',
            "SELECT id, question FROM meta.sql_exemplar"
            " WHERE status = 'enabled' AND embedding IS NULL AND ds_id = %s",   # 同上：meta::DS_PRED
            "UPDATE meta.sql_exemplar SET embedding = %s WHERE id = %s AND ds_id = %s",
            (ds,), (ds,), True)
        # 元素注册表向量（SuperSonic SchemaMapper 元素召回；search_text 变了自动 NULL 待重建）
        # 建表 DDL 的事实源在 Rust `meta::bootstrap_meta`；这份兜底副本必须跟着带 ds_id，
        # 否则「先跑 build 再启服务」会建出没有 ds_id 的表，下面那条 SQL 当场报缺列。
        cur.execute("CREATE TABLE IF NOT EXISTS meta.element("
                    "element_id text PRIMARY KEY, kind text NOT NULL, name text NOT NULL,"
                    "aliases text[] NOT NULL DEFAULT '{}', ref_expr text NOT NULL DEFAULT '',"
                    "description text NOT NULL DEFAULT '', search_text text NOT NULL DEFAULT '',"
                    "status text NOT NULL DEFAULT 'active',"
                    "ds_id text NOT NULL DEFAULT 'dms')")
        cur.execute(f"ALTER TABLE meta.element ADD COLUMN IF NOT EXISTS embedding vector({DIM})")
        n_el = _revec(
            cur, '元素',
            "SELECT element_id, search_text FROM meta.element"
            " WHERE status = 'active' AND embedding IS NULL AND ds_id = %s",    # 同上：meta::DS_PRED
            "UPDATE meta.element SET embedding = %s WHERE element_id = %s AND ds_id = %s",
            (ds,), (ds,))
        cur.execute("DROP INDEX IF EXISTS meta.idx_element_hnsw")
        cur.execute("CREATE INDEX idx_element_hnsw ON meta.element"
                    " USING hnsw (embedding vector_cosine_ops)")
        n_ds = _revec_datasources(cur)
    finally:
        pg.close()                  # 中途异常也不许残留连接（revec 同口径）
    print(f'完成[ds={ds}]：{n_tbl} 表 / {n_ex} 语料问句 / {n_el} 元素 / {n_ds} 数据源'
          f' 向量化 + HNSW 索引', flush=True)

def _revec_datasources(cur):
    """向量选源（`pipeline::select_source` → `meta::nearest_datasources`）的唯一写入点。
    ⚠️ 这里**不加也不能加 ds 限定**：meta.datasource 是 ds 注册表本身（Rust 那条漂移守卫
       也把它列为豁免），按 ds 过滤就只有当前源有向量 → 选源永远选不到别的源。
    只处理 embedding IS NULL：Rust `upsert_datasource` 在 description 变更时置 NULL 作失效
    （指 semantic/registry/datasource.rs 的主 upsert；它还有第二条 `register_upload_datasource`
    路径改 description **不清** embedding —— 陈旧向量的真正修复在 Rust 侧那条，这里口径不变）。
    文本 = name + description，与 `pick_by_llm` 给模型看的两个字段一致；
    问句侧是 embed_query（带指令前缀），故这里是文档侧 is_query=False，同 table_doc。"""
    return _revec(cur, '数据源',
                  "SELECT ds_id, name || '。' || description FROM meta.datasource"
                  " WHERE status = 'active' AND embedding IS NULL",
                  "UPDATE meta.datasource SET embedding = %s WHERE ds_id = %s")

# ============ 第五个 build 目标：知识库向量补齐（revec）============
# 缺陷背景（已实测）：`knowledge/src/ingest.rs` 在向量服务不可用时把文档停在 `chunked` 并写
# 「向量服务不可用，稍后可重建」——**而重建它的实现者一直不存在**。
# 后果不是「少一路召回」而是那份文档基本永久检索不到：`knowledge/src/retrieve.rs` 里
# 中文 FTS 实测恒 0（该路已换成单号/型号精确匹配，只保 ASCII token）、trgm 又有阈值挡着。
# 这两句不是推测，是量过的：14 个块 × 5 道题 = 70 组，**FTS 命中恒 0**；
# 「一线城市出差住宿费上限是多少」对答案所在块的 word_similarity 只有 0.267（<0.3，被挡）。
# 于是向量为空 = 零命中，而 UI 上那份文档显示「已入库」。
# 实测复现（报销制度.md，5 块）：embedding 置 NULL + status 退回 chunked → 同一问句
# citations **0 条**、回答「知识库里没有相关内容。」；跑本节后 citations 4 条、答出「500 元/晚」，
# 且 5 个块的 embedding md5 与置 NULL 前**逐个相同**（证明这里算的是文档侧向量、
# 和当初 ingest 同一个空间 —— 库因此也真回到了原状）。
#
# 一批的块数。实测本机：64 块（合 9675 字）一次 embed 2.21s，模型加载另 0.57s。
# 一批的文本总量因此是 ~10KB 量级 —— 几万块的库也不会被整份拉进内存（这是分批的全部理由；
# 每批立即回写 + autocommit，中途挂掉已补的不回滚，重跑从游标之后接着补）。
KB_BATCH = 64
KB_RECIPE = 1

def kb_embedding_text(name, folder_path, heading_path, body):
    """recipe v1；必须与 knowledge::store::chunk_embedding_text 逐字一致。"""
    return (f'文件：{name or ""}\n目录：{folder_path or "/"}\n章节：'
            f'{heading_path or "正文"}\n\n{body or ""}')

# 键集游标（`chunk_id > %s`）而不是「反复重查 embedding IS NULL」：失败的批**保持 NULL**，
# 重查同一条件就是死循环（永远拿回同一批）。游标只前进，失败的行落进「仍缺 K」。
KB_SEL = ("SELECT c.chunk_id,c.doc_id,d.name,c.folder_path,c.heading_path,c.text,"
          " c.embedding_text,c.embedding_recipe FROM kb.chunk c JOIN kb.doc d ON d.doc_id=c.doc_id"
          " WHERE c.embedding IS NULL AND c.embedding_recipe=%s AND c.chunk_id>%s"
          " AND d.status='chunked' AND d.enabled=true"
          " AND (d.effective_from IS NULL OR d.effective_from<=CURRENT_DATE)"
          " AND (d.effective_to IS NULL OR d.effective_to>=CURRENT_DATE)"
          " ORDER BY chunk_id LIMIT %s")
# `::vector` 显式转换：驱动对 str 参数的推断类型不一样（psycopg2 送 unknown 能隐式转，
# 送成 text 的驱动会当场报「column is of type vector but expression is of type text」）。
KB_UPD = ("UPDATE kb.chunk c SET embedding_text=%s,embedding_recipe=%s,embedding=%s::vector"
          " FROM kb.doc d WHERE c.chunk_id=%s AND c.embedding_text=%s AND c.embedding_recipe=%s"
          " AND c.embedding IS NULL AND d.doc_id=c.doc_id AND d.status='chunked' AND d.enabled=true"
          " AND (d.effective_from IS NULL OR d.effective_from<=CURRENT_DATE)"
          " AND (d.effective_to IS NULL OR d.effective_to>=CURRENT_DATE)")
KB_MISS = ("SELECT count(*) FROM kb.chunk c JOIN kb.doc d ON d.doc_id=c.doc_id"
           " WHERE (c.embedding IS NULL OR c.embedding_recipe<>%s) AND d.status='chunked' AND d.enabled=true"
           " AND (d.effective_from IS NULL OR d.effective_from<=CURRENT_DATE)"
           " AND (d.effective_to IS NULL OR d.effective_to>=CURRENT_DATE)")
KB_DOCS = ("SELECT doc_id,count(*),count(*) FILTER (WHERE embedding IS NULL OR embedding_recipe<>%s)"
           " FROM kb.chunk WHERE doc_id = ANY(%s) GROUP BY doc_id")
# `AND status = 'chunked'` 留在 SQL 里：failed / 已 embedded 的行不许被本节碰。
# error 只清那一句降级文案（其它文案可能是「表格已入知识库，建表失败」——那是真问题，不许抹）。
# 文案是 ingest.rs 的字面量，这里对不上时**只会不清**（不会误清），刻意选可失败安全的那一侧。
KB_PROMOTE = ("UPDATE kb.doc SET status = 'embedded',"
              " error = CASE WHEN error = %s THEN '' ELSE error END, updated_at = now()"
              " WHERE doc_id = ANY(%s) AND status = 'chunked' AND enabled=true"
              " AND (effective_from IS NULL OR effective_from<=CURRENT_DATE)"
              " AND (effective_to IS NULL OR effective_to>=CURRENT_DATE)")
DOWNGRADE_MSG = '向量服务不可用，稍后可重建'


def _promotable(n_chunks, n_missing):
    """能否把这份文档从 `chunked` 推到 `embedded`：**有块 且 一块不缺**。

    `n_chunks > 0` 不是废话：0 块的文档推成 embedded 就是「UI 显示已入库、其实一个字
    都检索不到」—— 正是本节要修的那个缺陷形态，不许在修它的代码里再造一遍。"""
    return n_chunks > 0 and n_missing == 0


def revec_exit(scanned, fixed, still):
    """退出码：**仍缺一行就非 0**。补了一半报成功正是这个缺陷活下来的方式
    （ingest 那边也是「记一句话然后返回 Ok」）。scanned/fixed 只进报告，不参与判定。"""
    return 1 if still else 0


def revec_chunks(cur, embed_fn=embed, batch=KB_BATCH):
    """扫空向量 → 按批补 → 按文档推进状态。返回 `(扫到, 补上, 仍缺, 推进的文档数)`。

    `embed_fn` 默认 `embed`（**文档侧** `is_query=False`，与 `EmbedClient::embed_passages`
    一致 —— 用 query 侧向量回写会让这些块和其它块不在同一个向量空间，检索时恒排在后面）。
    参数化只为 selftest：「一批向量化失败」这条路径不连库也得能验。
    **不截断正文**（`_revec` 那边截 1000 字是为注册表长描述）：ingest 时 Rust 侧原文入模型，
    这里截了就会算出和当初不同的向量；块本来 ≤480 token，模型 512 窗口自己会管。"""
    scanned = fixed = 0
    last, docs = 0, set()
    while True:
        cur.execute(KB_SEL, (KB_RECIPE, last, batch))
        rows = cur.fetchall()
        if not rows:
            break
        last = rows[-1][0]          # 游标先前进：失败的批也不许重来（否则死循环）
        scanned += len(rows)
        try:
            texts = [kb_embedding_text(name, folder, heading, body)
                     for _, _, name, folder, heading, body, _, _ in rows]
            vecs = embed_fn(texts)
            if len(vecs) != len(texts):
                raise RuntimeError(f'向量服务返回 {len(vecs)}/{len(texts)} 条，条数不符按本批失败处理')
        except Exception as e:
            # 保持 NULL 并让它计入「仍缺」——静默推进状态才是本节要修的缺陷
            print(f'  chunk …{last} 这批向量化失败，保持 NULL：{e}', flush=True)
            continue
        for (cid, doc, name, folder, heading, body, old_text, old_recipe), text, v in zip(rows, texts, vecs):
            cur.execute(KB_UPD, (text, KB_RECIPE, _vlit(v), cid, old_text, old_recipe))
            if cur.rowcount:
                fixed += 1
                docs.add(doc)
        print(f'  已补到 chunk {last}（{fixed}/{scanned}）', flush=True)
    promoted = _promote(cur, sorted(docs)) if docs else 0
    cur.execute(KB_MISS, (KB_RECIPE,))
    # 以库里的实际计数为准，不是 scanned - fixed：UPDATE 影响 0 行、别人并发置 NULL 都要算进来。
    # 代价（刻意选的）：正好有一次上传在「块已落、向量未回写」的窗口里时它会被算进 K → 非 0 退出。
    # 宁可多喊一次，也不要「补一半报成功」——后者正是这个缺陷的成因。
    still = cur.fetchone()[0]
    return scanned, fixed, still, promoted


def _promote(cur, docs):
    """这些文档里「一块不缺」的 → `embedded`，返回真正推进的份数。
    判据用**库里现在的计数**而不是「本次补了它几块」：并发上传/另一次 revec 都可能改动，
    只有 count 是事实源（也让重跑天然幂等 —— 已是 embedded 的被 SQL 里那条 status 挡住）。"""
    cur.execute(KB_DOCS, (KB_RECIPE, docs))
    ok = [d for d, n, miss in cur.fetchall() if _promotable(n, miss)]
    if not ok:
        return 0
    cur.execute(KB_PROMOTE, (DOWNGRADE_MSG, ok))
    return cur.rowcount


def revec():
    """`revec` / `build --revec` 的入口。返回退出码（仍缺即非 0）。"""
    import psycopg2
    pg = psycopg2.connect(connect_timeout=5, **pg_conf())
    pg.autocommit = True            # 每批即时落地：中途挂掉也不回滚已补好的（重跑接着补）
    cur = pg.cursor()
    # 卡在锁上不许无声等下去：本节是运维手动跑的，超时要响亮。
    # 口径注意：这两个只管**单条 SQL**，最慢的 embed 根本不是 SQL（一批 64 块实测 2.21s，
    # 不受 statement_timeout 约束）。实测 6 行那趟一共 10 条 SQL + 模型加载 + embed 合计 0.97s，
    # 所以 60s 对单条语句是纯兜底：真触发说明有锁或有全表扫，那时候该报错而不是等。
    # 两条 SET 分两次执行：塞在一个 execute 里是靠 psycopg2 简单查询协议才生效的，读着像会报错
    cur.execute("SET statement_timeout = '60s'")
    cur.execute("SET lock_timeout = '5s'")
    try:
        scanned, fixed, still, promoted = revec_chunks(cur)
    finally:
        pg.close()
    print(f'revec：扫到 {scanned} 行缺向量 / 补上 {fixed} 行 / 仍缺 {still} 行 / '
          f'状态推进 {promoted} 份文档', flush=True)
    return revec_exit(scanned, fixed, still)


# /embed 一次请求的条数上限：Rust 侧 embed.rs BATCH=64，这里给 4 倍余量；
# 不封顶的批量会长时间占住 `_EMBED_LOCK`（推理串行），把其它请求全饿死
EMBED_MAX_TEXTS = 256

def _int_param(body, name, default):
    """缺省/JSON null 用默认；显式给值必须能转 int（'abc' → 400 不是 500）。
    🔴 不许写 `body.get(name) or default`：显式 0 会被吞成默认（overlap=0 正是「关重叠」）。"""
    v = body.get(name)
    if v is None:
        return default
    try:
        return int(v)
    except (TypeError, ValueError):
        raise ParseError('bad_request', f'{name} 必须是整数：{v!r}', 400)

def handle_post(path, body):
    """POST 路由。/parse、/chunk **精确匹配**（/parseXYZ 不许进 /parse），未知路径按 /embed
    处理（兼容原来忽略 path 的行为）。字段类型错是 400（调用方的错）不是 500（我们的错）：
    path 传 list、blocks 传 dict、texts 传 "abc" 以前是 TypeError/AttributeError 500，
    而 texts 传 "abc" 更糟 —— 按字符逐一向量化，静默错。"""
    path = path.split('?')[0]          # 查询串不影响路由
    if path == '/parse':
        p, m = body.get('path'), body.get('mime')
        if (p is not None and not isinstance(p, str)) or (m is not None and not isinstance(m, str)):
            raise ParseError('bad_request', 'path/mime 必须是字符串', 400)
        return parse_doc(p or '', m or '')
    if path == '/chunk':
        blocks = body.get('blocks')
        if blocks is None:
            blocks = []
        if not isinstance(blocks, list) or any(not isinstance(b, dict) for b in blocks):
            raise ParseError('bad_request', 'blocks 必须是对象数组', 400)
        return {'chunks': chunk_blocks(blocks, _int_param(body, 'target_tokens', TARGET_TOKENS),
                                       _int_param(body, 'overlap', OVERLAP))}
    texts = body.get('texts')
    if texts is None:
        texts = []
    if not isinstance(texts, list) or any(not isinstance(t, str) for t in texts):
        raise ParseError('bad_request', 'texts 必须是字符串数组', 400)
    if len(texts) > EMBED_MAX_TEXTS:
        raise ParseError('too_large', f'texts {len(texts)} 条超上限 {EMBED_MAX_TEXTS}', 413)
    # query 缺省 False（文档侧）：漏传就静默进 query 向量空间是错的方向（与 passages 不同空间，
    # 检索恒排后）。Rust 端（embed.rs build_body）与 tools/ 探针都显式传 query，不受缺省影响。
    return {'embeddings': embed(texts, is_query=bool(body.get('query', False))) if texts else []}

# 请求体上限：/chunk 的 blocks 是大头（整份文档的块，实测量级 ~1MB）；32MB 是防内存 DoS
# 的兜底，不是业务限制 —— 超过直接 413，不许全读进内存（本进程还托着 /embed 的模型）
MAX_BODY_BYTES = 32 * 1024 * 1024

def _content_length(v):
    """Content-Length 解析：畸形/负值按 400 拒（以前 `int()` 在 try 外，
    ValueError 直接掐断连接、客户端什么响应都拿不到）。"""
    try:
        n = int(v or 0)
    except ValueError:
        raise ParseError('bad_request', f'Content-Length 非法：{v!r}', 400)
    if n < 0:
        raise ParseError('bad_request', f'Content-Length 非法：{v!r}', 400)
    return n

def serve(port=8077, host='127.0.0.1'):
    # host 显式可选（默认回环不松）：Linux 服务器部署时容器要经 docker 网桥（172.17.0.1）
    # 访问本服务 —— 绑网桥地址即可，0.0.0.0 会把解析/向量面暴露给公网。
    from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
    embedder()
    class H(BaseHTTPRequestHandler):
        def log_message(self, *a): pass
        def do_GET(self):
            # 健康检查（run.ps1 常驻化轮询用）
            if self.path.split('?')[0] == '/health':   # 探活常带 ?ts=… 防缓存
                resp = json.dumps({'ok': True, 'model': MODEL, 'dim': DIM,
                                   'parse_ok': parse_ok(), 'parse_caps': parse_caps()},
                                  ensure_ascii=False).encode()
                self.send_response(200)
            else:
                resp = b'{"error":"not found"}'
                self.send_response(404)
            self.send_header('Content-Type', 'application/json')
            self.send_header('Content-Length', str(len(resp)))
            self.end_headers()
            self.wfile.write(resp)
        def do_POST(self):
            try:
                n = _content_length(self.headers.get('Content-Length'))
                if n > MAX_BODY_BYTES:
                    raise ParseError('too_large', f'请求体 {n} 字节超上限 {MAX_BODY_BYTES}', 413)
                try:
                    body = json.loads(self.rfile.read(n) or b'{}')
                except json.JSONDecodeError as e:
                    raise ParseError('bad_json', f'请求体不是合法 JSON：{e}', 400)
                if not isinstance(body, dict):
                    raise ParseError('bad_request', '请求体必须是 JSON 对象', 400)
                resp = json.dumps(handle_post(self.path, body), ensure_ascii=False).encode()
                self.send_response(200)
            except ParseError as e:
                resp = json.dumps(e.payload, ensure_ascii=False).encode()
                self.send_response(e.status)
            except Exception as e:
                # 细节只进服务端日志：响应体回 str(e) 会把服务器绝对路径/内部形状泄给客户端
                # （serve 支持绑非回环地址，那是信息外泄面）
                print(f'[500] {self.path}: {type(e).__name__}: {e}', flush=True)
                resp = json.dumps({'error': 'internal'}, ensure_ascii=False).encode()
                self.send_response(500)
            self.send_header('Content-Type', 'application/json; charset=utf-8')
            self.send_header('Content-Length', str(len(resp)))
            self.end_headers()
            self.wfile.write(resp)
    # 🔴 ThreadingHTTPServer，不是 HTTPServer：stdlib 的 `socketserver.TCPServer` **串行**处理，
    # 而 /parse（上限 120s）/chunk /embed 全在同一个 handler 里。于是一次文档上传期间，
    # 每个问句的 embed 请求排在它后面 → 各自撞客户端 3s 超时
    # （`connector/src/embed.rs` TIMEOUT_SECS）→ send 失败把**进程级共享**的熔断打到 300s
    # → 之后 5 分钟语义缓存整条不生效、KB 向量路降级、表召回向量路降级，**全部静默**。
    # /health 探针也一起被堵（容器编排读成「服务挂了」）。
    # 同一件事 `docker/parser/parse_service.py:142` 早就修了，宿主这份漏了 —— 一份实现两个入口，
    # 只修一个入口等于没修。判据：`_selftest_serve_unblocked`（换回 HTTPServer 立刻红）。
    # 模型推理仍然串行（`_EMBED_LOCK`），并发的只有 IO 与解析。
    srv = ThreadingHTTPServer((host, port), H)
    # 先 bind 成功再报就绪：端口被占时上面那行就抛 OSError —— 先说就绪再炸会误导运维。
    # 启动就把「哪些扩展名不可用 + 为什么」打全：这条日志是运维唯一会看的地方，
    # 只打一个 True/False 的字典等于让人自己去猜缺哪个包。
    bad = [f"{e}（{c['why']}）" for e, c in sorted(parse_caps().items()) if not c['ok']]
    print(f'embed 服务就绪 :{port}（{MODEL}, {DIM}维）解析能力 {parse_ok()}'
          + (''.join(f'\n  ⛔ {b}' for b in bad) if bad else ''), flush=True)
    srv.serve_forever()

def selftest():
    """parse+chunk 自检：只用标准库（md/csv/json/html 四类），不依赖任何解析库与模型。
    末尾把**每种扩展名**的可用/不可用列全 —— 本机（Smart App Control 拦编译产物）预期
    pdf/docx/pptx/图片全是「不可用」，**那是正确输出**：判据是「不可用时明确说不可用」，
    不是「本机必须全可用」。真解析质量在容器里验（`tools/parse_probe.py`）。"""
    d = tempfile.mkdtemp()
    md = os.path.join(d, 'a.md')
    with open(md, 'w', encoding='utf-8') as f:
        f.write('# 一级标题\n开头一段。\n\n## 二级标题\n'
                + ''.join(f'第{i}句正文内容。' for i in range(200)))
    r = parse_doc(md)
    assert [b['heading_path'] for b in r['blocks']] == \
        ['一级标题', '一级标题', '一级标题 > 二级标题', '一级标题 > 二级标题'], r['blocks']
    ch = chunk_blocks(r['blocks'])
    assert len(ch) > 2 and all(c['tokens'] <= MAX_TOKENS for c in ch), ch
    assert ch[0]['heading_path'] == '一级标题' and ch[-1]['heading_path'] == '一级标题 > 二级标题'
    oc = int(OVERLAP * CHARS_PER_TOKEN)
    assert ch[2]['text'].startswith(ch[1]['text'][-oc:]), '同段内块间重叠丢了'
    cv = os.path.join(d, 'b.csv')
    with open(cv, 'w', encoding='utf-8') as f:
        f.write('客户,金额\n甲,1\n乙,2\n')
    r2 = parse_doc(cv)
    assert r2 == {'blocks': [], 'page_count': 0,
                  'sheets': [{'name': 'b', 'header': ['客户', '金额'],
                              'rows': [['甲', '1'], ['乙', '2']]}]}, r2
    # json：合法 → ```json 美化代码块；非法 → 原文入库且 notes 留痕（不许静默降级）
    js = os.path.join(d, 'c.json')
    with open(js, 'w', encoding='utf-8') as f:
        f.write('{"name": "甲", "tags": ["a", "b"]}')
    r3 = parse_doc(js)
    assert r3['blocks'][0]['text'].startswith('```json'), r3
    assert '"name": "甲"' in r3['blocks'][0]['text'], r3
    bad = os.path.join(d, 'bad.json')
    with open(bad, 'w', encoding='utf-8') as f:
        f.write('{not json')
    r4 = parse_doc(bad)
    assert r4['blocks'] and r4.get('notes') and 'JSON 校验失败' in r4['notes'][0], r4
    # html：去标签留文本；script/style 整棵丢弃；实体反转义
    ht = os.path.join(d, 'e.html')
    with open(ht, 'w', encoding='utf-8') as f:
        f.write('<html><head><title>制度</title><style>.a{color:red}</style></head>'
                '<body><h1>报销</h1><p>上限 &amp; 3000</p><script>alert(1)</script></body></html>')
    r5 = parse_doc(ht)
    joined = '\n'.join(b['text'] for b in r5['blocks'])
    assert '上限 & 3000' in joined and '报销' in joined, r5
    assert 'alert' not in joined and 'color' not in joined and '<p>' not in joined, r5
    open(os.path.join(d, 'x.rar'), 'wb').close()      # 真的没登记的扩展名
    for path, err in ((os.path.join(d, 'x.rar'), 'unsupported'), (os.path.join(d, 'nope.md'), 'not_found')):
        try:
            parse_doc(path)
            raise AssertionError(f'{path} 应该报 {err}')
        except ParseError as e:
            assert e.payload['error'] == err, e.payload
    caps = _selftest_caps(d)
    _selftest_pages()
    _selftest_pdf_scan(d)
    _selftest_xlsx_dims(d)
    _selftest_table_chunks()
    _selftest_revec()
    _selftest_md_heading()
    _selftest_text_edges(d)
    _selftest_pdf_runtime_fallback()
    _selftest_pptx_title(d)
    _selftest_image_frames(d)
    _selftest_emit_guard()
    _selftest_revec_len_guard()
    _selftest_handle_post()
    dt = _selftest_serve_unblocked()
    _selftest_http_errors()
    shutil.rmtree(d, ignore_errors=True)     # 临时目录不留（以前每跑一次 selftest 留一个）
    print(f'selftest ok: md块={len(r["blocks"])} 分块={len(ch)} tokens={[c["tokens"] for c in ch]}'
          f'（含 revec 纯逻辑 + 能力表双向一致 + /embed 忙 2s 时 /health {dt * 1000:.0f}ms 返回）'
          f'\n  parse_ok={parse_ok()}', flush=True)
    for ext, c in sorted(caps.items()):
        print(f"  {'✅' if c['ok'] else '⛔'} {ext:<10}{c['why']}", flush=True)


def _selftest_pdf_scan(tmpdir):
    """扫描件「OCR 档」（Y1）的钉。四段，每段对应一种不许回去的形态：
    ① 阈值纯逻辑：页均 <50 或全文 <200 判扫描件（裁决阈值，环境变量可覆盖 —— 改默认先改钉）
    ② 补页编排：桩 OCR 驱动 `_pdf_ocr_fill` —— 成功补页/单页失败留痕/整份响亮失败/页数护栏
    ③ 真实夹具：造一份**无文本** PDF，断言降级链不炸（确定性 ParseError，不是 500/静默返空）；
       fitz 环境再造一份**有墨图像页** PDF，真跑一次 OCR 档
    ④ 两档自报：`parse_caps()['.pdf']` 必须带 text|ocr 两档，且与真实可用性同源"""
    # ① 阈值钉
    assert _pdf_is_scanned(0, 3) and _pdf_is_scanned(199, 3)          # 全文 < 200
    assert _pdf_is_scanned(400, 10)                                   # 页均 40 < 50（全文 ≥200 也算）
    assert not _pdf_is_scanned(200, 3) and not _pdf_is_scanned(600, 3)
    assert not _pdf_is_scanned(0, 0), '0 页退化输入不许误判扫描件'
    assert _pdf_low_text_pages({1: 0, 2: 49, 3: 50, 4: 60}, 4) == [1, 2]
    assert _pdf_low_text_pages({}, 2) == [1, 2]                       # 缺席页按 0 计
    assert _page_chars(' 第 一 章　') == 3                            # 口径：非空白字符（含全角空格）
    # ② 补页编排（桩 OCR，不碰文件/引擎）
    ok = lambda p, i: [{'text': f'第{i}页内容' * 20, 'page': i, 'heading_path': f'第 {i} 页（OCR）'}]
    def bad(p, i):
        raise ParseError('unsupported', 'OCR 引擎缺席')
    empty = lambda p, i: []
    # 混合件：第 2 页低文本 → OCR 补上、按页号排回、notes 留痕
    out = [{'text': '正' * 400, 'page': 1, 'heading_path': ''}]
    notes = _pdf_ocr_fill('x.pdf', 2, {1: 400, 2: 0}, out, ocr_fn=ok)
    assert [b['page'] for b in out] == [1, 2] and notes and 'OCR 补' in notes[0], (out, notes)
    # 单页 OCR 失败不炸整份：原文保留 + 留痕
    out = [{'text': '正' * 400, 'page': 1, 'heading_path': ''}]
    notes = _pdf_ocr_fill('x.pdf', 2, {1: 400, 2: 0}, out, ocr_fn=bad)
    assert len(out) == 1 and notes and 'OCR 未成' in notes[-1], (out, notes)
    # 整份扫描件 + OCR 一字未补 → no_text_layer（不许带零星字符静默入库）；
    # 全文 <200 的「低文本量真文本」（3×60=180）同口径 —— 裁决阈值钉死，不许松
    for chars in ({1: 0, 2: 10}, {1: 60, 2: 60, 3: 60}):
        out = [{'text': '残' * c, 'page': i, 'heading_path': ''} for i, c in chars.items() if c]
        try:
            _pdf_ocr_fill('x.pdf', len(chars), chars, out, ocr_fn=bad)
            raise AssertionError(f'{chars} 应判扫描件报 no_text_layer')
        except ParseError as e:
            assert e.payload['error'] == 'no_text_layer', e.payload
    # 页均腿判扫描但全文 ≥200：OCR 全败时**保留**文本层 + 留痕（不升级为整份失败）
    chars = {i: 40 for i in range(1, 11)}
    out = [{'text': '残' * 40, 'page': i, 'heading_path': ''} for i in chars]
    notes = _pdf_ocr_fill('x.pdf', 10, chars, out, ocr_fn=bad)
    assert len(out) == 10 and notes and 'OCR 未成' in notes[-1], (len(out), notes)
    # 页数护栏：低文本页 > cap → too_large（不「OCR 前 N 页然后报已入库」）
    try:
        _pdf_ocr_fill('x.pdf', 5, {i: 0 for i in range(1, 6)}, [], ocr_fn=ok, cap=2)
        raise AssertionError('超 cap 应报 too_large')
    except ParseError as e:
        assert e.payload['error'] == 'too_large', e.payload
    # ocr_fn 返空按失败记（引擎读不出字不算「已补」），整份扫描件同样响亮失败
    try:
        _pdf_ocr_fill('x.pdf', 1, {1: 0}, [], ocr_fn=empty)
        raise AssertionError('OCR 返空应报 no_text_layer')
    except ParseError as e:
        assert e.payload['error'] == 'no_text_layer', e.payload
    # ③ 真实夹具：无文本 PDF（pypdf/fitz 有哪个用哪个；都没有 = pdf 能力本身不可用，跳过）
    fixture = os.path.join(tmpdir, 'scan_blank.pdf')
    made = False
    if _have('pypdf'):
        from pypdf import PdfWriter
        w = PdfWriter()
        w.add_blank_page(612, 792)
        w.add_blank_page(612, 792)
        with open(fixture, 'wb') as f:
            w.write(f)
        made = True
    elif _have('fitz'):
        import fitz
        d0 = fitz.open()
        d0.new_page()
        d0.new_page()
        d0.save(fixture)
        d0.close()
        made = True
    if made:
        try:
            r = parse_doc(fixture)
            # 空白页连 OCR 也读不出字 —— 能 200 只有一种解释：引擎把空白「读」出了内容。
            # 那时至少必须走了 OCR 档且留痕（不许返空块 —— 静默返空族）
            assert ''.join(b['text'] for b in r['blocks']).strip() and r.get('notes'), r
        except ParseError as e:
            # 本机形态（pypdf 兜底 + 无 OCR 引擎）：确定性失败 —— 不炸、不静默
            assert e.payload['error'] in ('no_text_layer', 'too_large'), e.payload
    # ③b 有墨的图像页夹具（仅 fitz 环境可造可验）：先排一版文字页，渲成 pixmap 再整页塞图，
    # 零第三方库造出「无文本层、图上有字」的真扫描件形态，真跑一次 OCR 档
    if _have('fitz'):
        import fitz
        src = fitz.open()
        tp = src.new_page()
        tp.insert_text((72, 120), 'DMS-SCAN-FIXTURE 7788', fontsize=28)
        pix = tp.get_pixmap(dpi=150)
        doc = fitz.open()
        ip = doc.new_page(width=tp.rect.width, height=tp.rect.height)
        ip.insert_image(ip.rect, pixmap=pix)
        img_pdf = os.path.join(tmpdir, 'scan_image.pdf')
        doc.save(img_pdf)
        doc.close()
        src.close()
        try:
            r = parse_doc(img_pdf)
            # OCR 引擎在：必须读出夹具暗号，且 notes 留痕「已用 OCR 补」
            assert any('7788' in b['text'] for b in r['blocks']), r
            assert any('OCR' in b.get('heading_path', '') for b in r['blocks']), r
            assert r.get('notes'), r
        except ParseError as e:
            # OCR 引擎缺席：no_text_layer 是确定性失败 —— 不炸、不静默返空
            assert e.payload['error'] == 'no_text_layer', e.payload
    # ④ 两档自报
    caps = parse_caps()
    if caps['.pdf']['ok']:
        t = caps['.pdf'].get('tiers')
        assert t and set(t) == {'text', 'ocr'}, caps['.pdf']
        assert t['text'] in ('pymupdf4llm', 'fitz', 'pypdf'), t
        assert t['ocr'] in ('qwen', 'tesseract', None), t
        if not _have('fitz'):
            assert t['ocr'] is None, 'fitz 缺席还上报 OCR 档可用 = 上报与真解析两套口径'
        assert 'OCR' in caps['.pdf']['why'], caps['.pdf']


def _selftest_xlsx_dims(tmpdir):
    """xlsx 静默丢列的钉（`_p_xlsx`）：WPS/ERP 导出常把 <dimension> 写小，read_only 轻信它
    就按声明截断。夹具 = 声明 A1:A1、实际 3 列的 xlsx，断言 3 列全回来。
    没有 openpyxl 时跳过（与 `_selftest_pdf_scan` 的真实夹具同纪律：能力缺席 ≠ 判据消失）。"""
    if not _have('openpyxl'):
        print('  ⏭️  xlsx 丢列判据跳过（openpyxl 不可用）', flush=True)
        return
    import openpyxl
    import zipfile
    src = os.path.join(tmpdir, 'dims_src.xlsx')
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = '台账'
    ws.append(['单号', '金额', '备注'])
    ws.append(['A-1', '10', '甲'])
    ws.append(['B-2', '20', '乙'])
    wb.save(src)
    bad = os.path.join(tmpdir, 'dims_bad.xlsx')
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(bad, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith('xl/worksheets/sheet'):
                # WPS/ERP 导出的实测形态：声明比实际小
                data = re.sub(rb'<dimension ref="[^"]*"/>', b'<dimension ref="A1:A1"/>', data)
            zout.writestr(item, data)
    r = parse_doc(bad)
    sheet = r['sheets'][0]
    assert sheet['header'] == ['单号', '金额', '备注'], f'静默丢列复现：{sheet}'
    assert sheet['rows'] == [['A-1', '10', '甲'], ['B-2', '20', '乙']], sheet


def _selftest_table_chunks():
    """表格行感知分块的钉（`_fill_table`）：每块必须重复表头、一行不丢不重、顺序不乱；
    同标题下的散文段仍走通用分块（互不染指）。"""
    rows = [f'| SO-2026-{i:03d} | {100 + i} 元 |' for i in range(60)]
    text = '# 一月台账\n\n| 单号 | 金额 |\n| --- | --- |\n' + '\n'.join(rows)
    ch = chunk_blocks([{'text': text, 'heading_path': '一月台账', 'page': None}],
                      target_tokens=60, overlap=0)
    assert len(ch) > 3, ch
    for c in ch:
        lines = c['text'].split('\n')
        assert lines[:4] == ['# 一月台账', '', '| 单号 | 金额 |', '| --- | --- |'], c
        assert all(l.startswith('| SO-') for l in lines[4:]), c
    got = [l for c in ch for l in c['text'].split('\n') if l.startswith('| SO-')]
    assert got == rows, '数据行丢/重/乱序'
    mix = chunk_blocks([{'text': '前文说明。' * 30, 'heading_path': 'H', 'page': None},
                        {'text': text, 'heading_path': 'H', 'page': None}],
                       target_tokens=60, overlap=0)
    assert any('前文说明' in c['text'] for c in mix)
    for c in mix:
        if c['text'].startswith('# 一月台账'):
            assert '前文说明' not in c['text'], c


def _selftest_caps(tmpdir):
    """能力表纪律。四条断言，每条都对应一种已经发生过的静默失败：
    ① 扩展名表与解析器表**双向一致** —— 「登记而不消费」（声明支持却没有实现者）与
       「实现了却没登记」（写了解析器但 /parse 永远走不到）本轮各抓到过一次
    ② MIME 表只许映到登记过的扩展名（映到没登记的 = 按 mime 上传时莫名 unsupported）
    ③ **上报不可用的扩展名，真去解析必须抛 `unsupported` 且带原因** —— 静默返空是本仓
       反复抓的失败族（status=embedded + chunk 0 + 界面「已入库」+ 问不出东西）
    ④ `parse_ok` 里 `pdf`/`xlsx` 两个键在：`parse_probe.py` 靠它们决定跳不跳，
       键改名它会静默全跳过并 exit 0（判据恒绿）"""
    impl = {n: f for n, f in globals().items() if n.startswith('_p_') and callable(f)}
    reg = {f for f, _, _ in CAPS.values()}
    assert set(impl.values()) == reg, (
        f'解析器登记不一致：写了却没登记 {sorted(n for n, f in impl.items() if f not in reg)}'
        f' / 登记的不是模块级 _p_* 函数 {sorted(e for e, (f, _, _) in CAPS.items() if f not in impl.values())}')
    assert set(MIME_EXT.values()) <= set(CAPS), sorted(set(MIME_EXT.values()) - set(CAPS))
    assert {'pdf', 'xlsx'} <= set(parse_ok()), parse_ok()
    caps = parse_caps()
    assert set(caps) == set(CAPS), '能力上报漏了扩展名'
    for ext, c in caps.items():
        if c['ok']:
            continue
        p = os.path.join(tmpdir, 'cap' + ext)   # 0 字节即可：门在依赖探测，轮不到读内容
        open(p, 'wb').close()
        try:
            got = parse_doc(p)
        except ParseError as e:
            assert e.payload['error'] == 'unsupported' and e.payload.get('detail'), (ext, e.payload)
            continue
        raise AssertionError(f'{ext} 上报不可用，解析却返回了 {got} —— 静默返空族')
    return caps


class _FakeCur:
    """selftest 的假游标：按 SQL **常量身份**（`is`）分派，不解析 SQL 文本。
    好处是 SQL 改了字它不会假装还在测同一条路 —— 认不出来的语句当场 AssertionError。"""

    def __init__(self, batches, groups, missing):
        self.batches, self.groups, self.missing = list(batches), groups, missing
        self.rows, self.cursors, self.updated, self.promoted = [], [], [], []
        self.rowcount = 0

    def execute(self, sql, params=()):
        self.rows = []
        if sql is KB_SEL:
            self.cursors.append(params[1])
            self.rows = self.batches.pop(0) if self.batches else []
        elif sql is KB_UPD:
            self.updated.append(params[3])
            self.rowcount = 1
        elif sql is KB_DOCS:
            self.rows = self.groups
        elif sql is KB_MISS:
            self.rows = [(self.missing,)]
        elif sql is KB_PROMOTE:
            self.promoted = list(params[1])
            self.rowcount = len(self.promoted)
        else:
            raise AssertionError(f'假游标不认识这条 SQL：{sql}')

    def fetchall(self):
        return self.rows

    def fetchone(self):
        return self.rows[0] if self.rows else None


def _selftest_revec():
    """revec 纯逻辑自检（不连库、不加载模型）。三条都是这个缺陷的实际形状：
    ① 一批向量化失败 → 那些行保持 NULL、计入「仍缺」，且**游标仍要前进**（不然死循环）
    ② 状态推进的条件是「有块且一块不缺」，不是「本次补过它」
    ③ 仍缺一行就非 0 退出（补一半报成功 = 复刻缺陷）"""
    def flaky(texts):
        if len(texts) == 1:            # 第二批只剩一行时假装向量服务抖了
            raise RuntimeError('模拟向量服务抖动')
        return [[0.5] * DIM] * len(texts)

    cur = _FakeCur(batches=[[
            (1, 'd1', 'a.md', '/制度', '甲', 'a', 'old-a', KB_RECIPE),
            (2, 'd1', 'a.md', '/制度', '乙', 'b', 'old-b', KB_RECIPE)],
            [(3, 'd2', 'b.md', '/', '', 'c', 'old-c', KB_RECIPE)]],
                   groups=[('d1', 2, 0), ('d2', 1, 1)], missing=1)
    got = revec_chunks(cur, flaky, batch=2)
    assert got == (3, 2, 1, 1), got                     # 扫 3 / 补 2 / 仍缺 1 / 推进 1
    assert cur.updated == [1, 2], cur.updated           # 失败那行没被回写
    assert cur.cursors == [0, 2, 3], cur.cursors        # 游标越过失败批：0→2→3 后取空退出
    assert cur.promoted == ['d1'], cur.promoted         # d2 还缺一块 → 不许推进
    assert revec_exit(*got[:3]) == 1, '仍缺 1 行还给 0 退出码 = 补一半报成功'
    assert revec_exit(0, 0, 0) == 0 and revec_exit(9, 9, 0) == 0
    assert not _promotable(0, 0), '0 块的文档不许推成 embedded'
    assert _promotable(3, 0) and not _promotable(3, 1)
    assert kb_embedding_text('制度.md', '/财务/报销', '第一章 > 范围', '正文') == \
        '文件：制度.md\n目录：/财务/报销\n章节：第一章 > 范围\n\n正文'

def _selftest_md_heading():
    """`_H_MD` 的钉：中文文档常写 `#一级标题`（# 后不空格），要认出层级；
    但 `#tag` / `#123` 不是标题（# 后首字符是 ASCII 字母数字的不放行）。"""
    b = md_blocks('#一级标题\n正文。\n\n##二级\n再一句。', None, [])
    assert [x['text'] for x in b] == ['一级标题', '正文。', '二级', '再一句。'], b
    assert b[-1]['heading_path'] == '一级标题 > 二级', b
    b = md_blocks('#tag\n#123', None, [])
    assert [x['heading_path'] for x in b] == [''], b      # 都不是标题，合成一段正文


def _selftest_text_edges(tmpdir):
    """文本族入口的边角钉（`_read_text`/`_p_text`/`parse_doc`）：
    ① UTF-16 BOM 先拦：Windows 导出的 txt 落入 gbk 会解成夹 NUL 的乱码而不是正确解码
    ② `---`/`***` 这种无词字符块不进 blocks（与 `_p_pdf` 同一过滤口径）
    ③ mime 匹配大小写不敏感（RFC 允许 `Text/Plain` 这种写法）"""
    p = os.path.join(tmpdir, 'u16.txt')
    with open(p, 'w', encoding='utf-16') as f:          # 带 BOM
        f.write('第二章\n带 BOM 的 UTF-16 内容。')
    r = parse_doc(p)
    joined = '\n'.join(b['text'] for b in r['blocks'])
    assert '带 BOM 的 UTF-16 内容。' in joined and '\x00' not in joined, r
    p = os.path.join(tmpdir, 'sep.md')
    with open(p, 'w', encoding='utf-8') as f:
        f.write('# 标题\n正文。\n\n---\n\n***\n\n下一段。')
    r = parse_doc(p)
    assert all(re.search(r'\w', b['text']) for b in r['blocks']), r['blocks']
    p = os.path.join(tmpdir, 'only_sep.md')
    with open(p, 'w', encoding='utf-8') as f:
        f.write('---\n\n***\n')
    assert parse_doc(p)['blocks'] == [], '分隔线-only 文档不该产出块'
    p = os.path.join(tmpdir, 'x.dat')                    # 没登记的扩展名，靠 mime 查表
    with open(p, 'w', encoding='utf-8') as f:
        f.write('纯文本内容一句。')
    r = parse_doc(p, 'Text/Plain; charset=UTF-8')
    assert r['blocks'], 'mime 大小写必须不影响查表'
    # `_cell` strip：xlsx/csv 里 ' 10 ' 这类带空白值不该原样进 sheets/表头
    p = os.path.join(tmpdir, 'pad.csv')
    with open(p, 'w', encoding='utf-8') as f:
        f.write('客户,金额\n甲, 10 \n')
    assert parse_doc(p)['sheets'][0]['rows'] == [['甲', '10']], '单元格带空白必须 strip'


def _selftest_pdf_runtime_fallback():
    """`_p_pdf` 的运行期降级钉：pymupdf4llm **import 成功**但 `to_markdown` 抛异常
    （损坏/异形 PDF）必须续降 `_pdf_fitz`，不许整份 500。桩模块驱动，不碰真 PDF。"""
    import types
    fake = types.ModuleType('pymupdf4llm')
    fake.to_markdown = lambda *a, **k: (_ for _ in ()).throw(RuntimeError('模拟损坏 PDF'))
    sentinel = ([{'text': 'x', 'page': 1, 'heading_path': ''}], 1, [], [])
    keep_mod, keep_fitz = sys.modules.get('pymupdf4llm'), globals()['_pdf_fitz']
    sys.modules['pymupdf4llm'] = fake
    globals()['_pdf_fitz'] = lambda path: sentinel
    try:
        assert _p_pdf('x.pdf') is sentinel, 'pymupdf4llm 运行期异常必须续降 _pdf_fitz'
    finally:
        if keep_mod is None:
            sys.modules.pop('pymupdf4llm', None)
        else:
            sys.modules['pymupdf4llm'] = keep_mod
        globals()['_pdf_fitz'] = keep_fitz


def _selftest_pptx_title(tmpdir):
    """pptx 标题去重的钉：标题 shape 也有 text_frame，不排除会同时进 heading_path 和正文，
    标题被向量化两遍。没有 python-pptx 时跳过（同 `_selftest_xlsx_dims` 纪律）。"""
    if not _have('pptx'):
        print('  ⏭️  pptx 标题判据跳过（python-pptx 不可用）', flush=True)
        return
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])   # 版式 0 带标题占位符
    slide.shapes.title.text = '报销制度'
    slide.placeholders[1].text = '上限三千。'
    p = os.path.join(tmpdir, 't.pptx')
    prs.save(p)
    r = parse_doc(p)
    assert len(r['blocks']) == 1, r['blocks']
    assert r['blocks'][0]['heading_path'] == '报销制度', r['blocks']
    assert '报销制度' not in r['blocks'][0]['text'] and '上限三千。' in r['blocks'][0]['text'], r['blocks']


def _selftest_image_frames(tmpdir):
    """多帧图片的三颗钉（桩 OCR 驱动，不碰引擎；PIL 缺席跳过）：
    ① 帧数超 OCR_PAGE_CAP → too_large（每帧一次千问 HTTP，N 帧会远超 Rust 120s 解析超时）
    ② 多帧 heading_path 带帧号（共用文件名会被 chunk 按 heading_path 合并出跨帧重复）；
       单帧图保持裸文件名
    ③ OCR 引擎抛的 ParseError 原样上抛，不许再包一层套娃 detail"""
    if not _have('PIL.Image'):
        print('  ⏭️  图片帧判据跳过（PIL 不可用）', flush=True)
        return
    from PIL import Image
    p = os.path.join(tmpdir, 'frames.gif')
    Image.new('RGB', (8, 8), 'white').save(
        p, save_all=True, append_images=[Image.new('RGB', (8, 8), 'black'),
                                        Image.new('RGB', (8, 8), 'red')])
    g = globals()
    keep = g['_ocr_qwen_frame'], g['_ocr_tesseract_frame'], g['OCR_PAGE_CAP']
    try:
        g['_ocr_qwen_frame'] = lambda frame: None
        g['_ocr_tesseract_frame'] = lambda frame: '第N帧文字'
        # ② 三帧 → heading_path 各不相同且带帧号；单帧图保持裸文件名
        blocks, frames, _, _ = _p_image(p)
        assert frames == 3 and len(blocks) == 3, (frames, blocks)
        assert [b['heading_path'] for b in blocks] == \
            ['frames.gif#f1', 'frames.gif#f2', 'frames.gif#f3'], blocks
        one = os.path.join(tmpdir, 'one.png')
        Image.new('RGB', (8, 8), 'white').save(one)
        b1 = _p_image(one)[0]
        assert [b['heading_path'] for b in b1] == ['one.png'], b1
        # ③ ParseError 原样上抛不套娃
        def _bad(frame):
            raise ParseError('unsupported', 'tesseract OCR 失败（lang=x）')
        g['_ocr_tesseract_frame'] = _bad
        try:
            _p_image(one)
            raise AssertionError('ParseError 应原样上抛')
        except ParseError as e:
            detail = e.payload.get('detail', '')
            assert 'tesseract OCR 失败' in detail and 'OCR 失败（tesseract' not in detail, e.payload
        # ① 帧数护栏：cap 调到 2，三帧 gif 响亮 too_large
        g['_ocr_tesseract_frame'] = keep[1]
        g['OCR_PAGE_CAP'] = 2
        try:
            _p_image(p)
            raise AssertionError('超帧数上限应报 too_large')
        except ParseError as e:
            assert e.payload['error'] == 'too_large' and '3' in e.payload['detail'], e.payload
    finally:
        g['_ocr_qwen_frame'], g['_ocr_tesseract_frame'], g['OCR_PAGE_CAP'] = keep


def _selftest_emit_guard():
    """`_emit` 的窗护栏钉：超 MAX_TOKENS 必须显式 raise（python -O 下 assert 会被剥掉，
    超窗块静默进库被 fastembed 截断）。"""
    try:
        _emit([], '超' * (MAX_TOKENS * 2), 'H', None)      # est_tokens ≈ 1.25×MAX_TOKENS
        raise AssertionError('超窗块必须 raise')
    except ValueError as e:
        assert '上限' in str(e), e


def _selftest_revec_len_guard():
    """向量条数与行数不符的钉（`_revec`/`revec_chunks`，与 Rust embed.rs
    「少返一行→整批 None」同一纪律）：`_revec` 响亮失败（zip 静默截断会把向量错位回写），
    `revec_chunks` 按「该批失败」处理（保持 NULL、计入仍缺、游标前进）。"""
    cur = _FakeCur(batches=[[(1, 'd1', 'a.md', '/', '', 'a', 'old-a', KB_RECIPE),
                             (2, 'd1', 'a.md', '/', '', 'b', 'old-b', KB_RECIPE)]],
                   groups=[], missing=2)
    got = revec_chunks(cur, lambda texts: [[0.5] * DIM] * (len(texts) - 1), batch=2)
    assert got == (2, 0, 2, 0), got                     # 少返一条 → 这批保持 NULL
    assert cur.updated == [] and cur.cursors == [0, 2], cur.cursors
    class _Cur:
        def execute(self, sql, params=()):
            self.rows = [(1, 't1'), (2, 't2')]
        def fetchall(self):
            return self.rows
    keep = globals()['embed']
    globals()['embed'] = lambda texts, is_query=False: [[0.5] * DIM] * (len(texts) - 1)
    try:
        try:
            _revec(_Cur(), '表', 'SEL', 'UPD')
            raise AssertionError('条数不符必须响亮失败')
        except RuntimeError as e:
            assert '不符' in str(e), e
    finally:
        globals()['embed'] = keep


def _selftest_handle_post():
    """HTTP 入口形状的钉（不起服务，直接打 `handle_post`）：
    ① 路由精确：/parseXYZ、/chunky 不许进对应 handler（未知路径按 /embed，兼容口径不变）
    ② 显式 `overlap: 0` 不许被吞成默认 60（`x or DEFAULT` 的坑）；缺省/None 仍用默认
    ③ 字段类型错 → 400：path 传 list / blocks 传 dict / texts 传 "abc" / overlap 传 'abc'
    ④ /embed 缺省 query=false（Rust 端总是显式传）；texts 条数封顶 413"""
    g = globals()
    keep = g['embed']
    seen = []
    g['embed'] = lambda texts, is_query=False: (seen.append(is_query), [[0.0] * DIM] * len(texts))[1]
    try:
        # ① 未知路径（含 /parseXYZ、/chunky）按 /embed；查询串不影响路由
        r = handle_post('/parseXYZ', {'texts': ['a']})
        assert 'embeddings' in r and seen[-1] is False, r
        r = handle_post('/chunky?x=1', {'texts': ['a'], 'query': True})
        assert 'embeddings' in r and seen[-1] is True, r
        # ④ 缺省 query = False
        handle_post('/embed', {'texts': ['a']})
        assert seen[-1] is False, seen
        # ② overlap=0 真的 0（与直接调 chunk_blocks 逐字节一致）；缺省/None 仍是默认 60
        blocks = [{'text': ''.join(chr(0x4e00 + i % 500) for i in range(300)),
                   'heading_path': 'H', 'page': 1}]
        want0 = chunk_blocks(blocks, 60, 0)
        assert handle_post('/chunk', {'blocks': blocks, 'target_tokens': 60, 'overlap': 0})['chunks'] \
            == want0, '显式 overlap=0 被吞成默认值'
        wantd = chunk_blocks(blocks, 60, OVERLAP)
        assert handle_post('/chunk', {'blocks': blocks, 'target_tokens': 60})['chunks'] == wantd
        assert handle_post('/chunk', {'blocks': blocks, 'target_tokens': 60, 'overlap': None})['chunks'] \
            == wantd, 'JSON null 应按缺省处理'
        assert want0 != wantd, '夹具没造出重叠差异，这条钉白搭'
        # /parse 正常路径仍通（文件不存在 → 404 not_found）
        try:
            handle_post('/parse', {'path': 'Z:/no/such/file.md'})
            raise AssertionError('应报 not_found')
        except ParseError as e:
            assert e.payload['error'] == 'not_found' and e.status == 404, e.payload
        # ③ 类型校验
        for path, body in (('/parse', {'path': ['x']}),
                           ('/parse', {'path': 'x.md', 'mime': 1}),
                           ('/chunk', {'blocks': {'text': 'x'}}),
                           ('/chunk', {'blocks': ['x']}),
                           ('/chunk', {'blocks': [], 'overlap': 'abc'}),
                           ('/embed', {'texts': 'abc'}),
                           ('/embed', {'texts': [1]})):
            try:
                handle_post(path, body)
                raise AssertionError(f'{path} {body} 应报 400')
            except ParseError as e:
                assert e.status == 400, (path, body, e.payload)
        # ④ 条数上限
        try:
            handle_post('/embed', {'texts': ['a'] * (EMBED_MAX_TEXTS + 1)})
            raise AssertionError('超上限应报 413')
        except ParseError as e:
            assert e.status == 413, e.payload
    finally:
        g['embed'] = keep


def _selftest_http_errors():
    """HTTP 错误形状的钉（真起一次 serve，桩 embed 抛错）：
    ① 非法 JSON → 400 bad_json（以前走 except Exception 回 500）
    ② 非对象 JSON（数组）→ 400 bad_request
    ③ handler 意外异常 → 500 但响应体**不带**异常原文（内部细节只进服务端日志）
    ④ /health 带查询串也认（探活防缓存写法 ?ts=…）
    ⑤ 超大/畸形 Content-Length → 413/400（以前超大全读进内存、畸形直接掐断连接）"""
    import socket
    import time
    g = globals()
    keep = g['embed'], g['embedder']
    def _boom(texts, is_query=False):
        raise RuntimeError('内部细节-D:/secret/path 不该外泄')
    g['embed'] = _boom
    g['embedder'] = lambda: None
    with socket.socket() as s:                 # 借一个空闲端口：别撞常驻的 8077
        s.bind(('127.0.0.1', 0))
        port = s.getsockname()[1]
    base = f'http://127.0.0.1:{port}'
    def post(raw):
        req = urllib.request.Request(base + '/embed', data=raw,
                                     headers={'Content-Type': 'application/json'})
        try:
            with urllib.request.urlopen(req, timeout=5) as r:
                return r.status, json.loads(r.read())
        except urllib.error.HTTPError as e:
            return e.code, json.loads(e.read())
    def raw_status(cl):                        # 手写 Content-Length 的原始请求
        with socket.socket() as c:
            c.settimeout(5)
            c.connect(('127.0.0.1', port))
            c.sendall(b'POST /embed HTTP/1.1\r\nHost: x\r\nContent-Length: ' + cl +
                      b'\r\nConnection: close\r\n\r\n')
            resp = b''
            while True:
                chunk = c.recv(4096)
                if not chunk:
                    break
                resp += chunk
        return resp.split(b'\r\n', 1)[0]
    try:
        threading.Thread(target=serve, args=(port,), daemon=True).start()  # 关不掉，同 unblocked 那条注释
        for _ in range(50):                    # 等监听就绪
            try:
                with urllib.request.urlopen(base + '/health?ts=1', timeout=5) as r:   # ④
                    assert r.status == 200 and json.loads(r.read())['ok']
                break
            except Exception:
                time.sleep(0.1)
        else:
            raise AssertionError(f'serve 没能在 5s 内起在 {port}')
        code, j = post(b'{not json')                                   # ①
        assert code == 400 and j['error'] == 'bad_json', (code, j)
        code, j = post(b'[1,2]')                                       # ②
        assert code == 400 and j['error'] == 'bad_request', (code, j)
        code, j = post(json.dumps({'texts': ['a']}).encode())          # ③
        assert code == 500 and '内部细节' not in json.dumps(j, ensure_ascii=False), (code, j)
        assert b' 413' in raw_status(b'999999999'), '超大 Content-Length 必须 413'      # ⑤
        assert b' 400' in raw_status(b'abc'), '畸形 Content-Length 必须 400'
    finally:
        g['embed'], g['embedder'] = keep

def _selftest_serve_unblocked():
    """**真起一次 serve**：/embed 慢的时候 /health 必须仍在 500ms 内返回。

    这条判据钉的是一个已经发生过的回归：同一件事在 `docker/parser/parse_service.py` 修过
    （ThreadingHTTPServer），宿主这份漏了。后果不是「慢」而是**静默错答** ——
    一次上传把每个问句的 embed 堵到 3s 超时，客户端进程级熔断 300s，
    之后 5 分钟语义缓存/KB 向量路/表召回向量路全降级且零报错。

    ⚠️ 为什么不写成「源码里有 ThreadingHTTPServer」那种断言：本仓刚被同一形态咬过一次 ——
    `check-arch.ps1` 的 reqwest 那条把**注释里的字**当真实用点。注释里提一句
    ThreadingHTTPServer、代码里照旧 `HTTPServer(...)`，字符串断言会绿。所以这里量的是墙钟。

    不加载模型（patch 掉 `embedder`/`embed`）：selftest 的纪律是不碰第三方库与 95MB 模型。"""
    import socket
    import time
    g = globals()
    keep = g['embed'], g['embedder']
    # 桩睡 2s。**别调小**：断言在 /embed 发出后 0.15s 才开始计时，剩余时间必须仍远超 500ms 阈值。
    # 第一版写 0.6s，单线程时剩 0.45s < 0.5s ⇒ 反向验证时这条判据是**绿的**（恒真）。
    # 修好后这 2s 不进 selftest 的墙钟：/health 立刻返回，embed 桩在 daemon 线程里被丢下。
    g['embed'] = lambda texts, is_query=False: (time.sleep(2), [[0.0] * DIM] * len(texts))[1]
    g['embedder'] = lambda: None
    with socket.socket() as s:                 # 借一个空闲端口：别撞常驻的 8077
        s.bind(('127.0.0.1', 0))
        port = s.getsockname()[1]
    base = f'http://127.0.0.1:{port}'

    def get_health(timeout=5):
        with urllib.request.urlopen(base + '/health', timeout=timeout) as r:
            return json.loads(r.read())

    try:
        # serve 不返回 server 对象（serve_forever 阻塞），这里起的实例关不掉 ——
        # 靠 daemon 线程 + 进程退出兜底；selftest 进程短命，可接受
        threading.Thread(target=serve, args=(port,), daemon=True).start()
        for _ in range(50):                    # 等监听就绪（bind 之后到 accept 之前有个窗口）
            try:
                assert get_health()['ok']
                break
            except Exception:
                time.sleep(0.1)
        else:
            raise AssertionError(f'serve 没能在 5s 内起在 {port}')
        threading.Thread(target=lambda: urllib.request.urlopen(
            urllib.request.Request(base + '/embed', json.dumps({'texts': ['a']}).encode(),
                                   {'Content-Type': 'application/json'}), timeout=10),
            daemon=True).start()
        time.sleep(0.15)                       # 让 /embed 先进 handler 睡着
        t0 = time.time()
        assert get_health()['ok']
        dt = time.time() - t0
        assert dt < 0.5, (f'/embed 睡 2s 期间 /health 等了 {dt:.3f}s —— serve 又变回单线程了；'
                          '一次上传会把每个问句的 embed 顶到超时 + 300s 熔断')
        return dt
    finally:
        g['embed'], g['embedder'] = keep


if __name__ == '__main__':
    import argparse
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument('mode', nargs='?', default='serve',
                    choices=('build', 'revec', 'serve', 'selftest'))
    ap.add_argument('port', nargs='?', type=int, default=8077, help='serve 端口（默认 8077）')
    ap.add_argument('host', nargs='?', default='127.0.0.1', help='serve 绑定地址（默认回环；容器跨网桥访问时给 172.17.0.1）')
    ap.add_argument('--ds', default='dms',
                    help='build 只处理该 ds_id 的注册表行（对应 Rust 侧 meta::DS_PRED，默认 dms）')
    ap.add_argument('--revec', action='store_true',
                    help='只跑第五个目标：补 kb.chunk 空向量 + 推进 kb.doc 状态（= mode revec）')
    a = ap.parse_args()
    if a.revec or a.mode == 'revec':
        sys.exit(revec())      # 仍缺即非 0：判定在 revec_exit，别在这里改成恒 0
    elif a.mode == 'build':
        build(a.ds)
    elif a.mode == 'selftest':
        selftest()
    else:
        serve(a.port, a.host)
